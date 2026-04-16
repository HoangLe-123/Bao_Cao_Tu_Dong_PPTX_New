import pandas as pd
import matplotlib.pyplot as plt
import numpy as np
# import matplotlib.dates as mdates
from pptx import Presentation
from pathlib import Path
from matplotlib.ticker import FixedLocator, FixedFormatter
from matplotlib.transforms import offset_copy
plt.rcParams["font.family"] = "Segoe UI"
plt.rcParams["axes.unicode_minus"] = False
# hoặc:
# plt.rcParams["font.family"] = "Yu Gothic"
# plt.rcParams["font.family"] = "Noto Sans CJK JP"
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from datetime import datetime

# ========== TUNABLES ==========
EXCEL_FILE = r"Hằng ngày_data.xlsx"
EXCEL_SHEET = "Hằng ngày"
EXCEL_YEAR_FILE = r"Tính sản xuất năm 2026.xlsx"
EXCEL_YEAR_SHEET = "NĂM 2026"
EXCEL_THUCTICH_FILE = r"ThucTich_Data"
EXCEL_THUCTICH_SHEET = "Sheet1"
EXCEL_COSTDOWN_FILE = r"Costdown nội tác 2026.xlsx"
EXCEL_COSTDOWN_SHEET = "Quản lý MT năm"
PPT_FILE = r"Bao_Cao_Hang_Ngay.pptx"
ROW_DATE = "Ngày làm việc"
OUT_DIR = Path("charts")
OUT_DIR.mkdir(exist_ok=True)

# Scale multipliers (tinh chỉnh để biểu đồ trông giống thủ công)
LEFT_Y_MULT = 1.3   # scale cho trục trái (tháng)
RIGHT_Y_MULT = 1.3  # scale cho trục phải (ngày)
MONTH_WIDTH_DAYS = 18
DAY_BAR_WIDTH = 0.6

# Debug: bật để in df_daily / df_month (False khi chạy chính thức)
DEBUG = False

# ========== CONFIGS ==========
GLOBAL_CONFIG = [
    {
        "name": "tong_ban_ve",
        "legend_prefix": "Bản vẽ",
        "rows": {
            "tt": "THỰC TÍCH SỐ BẢN VẼ HOÀN THÀNH TRONG NGÀY",
            "mt": "MỤC TIÊU SỐ BẢN VẼ HOÀN THÀNH TRONG NGÀY",
            "lk": "MỤC TIÊU LŨY KẾ BẢN VẼ HT TRONG NGÀY ( CÒN LẠI)",
        },
        "title": "TỔNG BẢN VẼ HOÀN THÀNH TRONG THÁNG ",
        "ppt_shape": "tong_ban_ve_shape"
    },
    
    {
        "name": "tong_so_pcs",
        "legend_prefix": "Số PCS",
        "rows": {
            "tt": "THỰC TÍCH SỐ PCS HOÀN THÀNH TRONG NGÀY",
            "mt": "MỤC TIÊU SỐ PCS HOÀN THÀNH TRONG NGÀY",
            "lk": "MỤC TIÊU LŨY KẾ SỐ PCS HT TRONG NGÀY ( CÒN LẠI)",
        },
        "title": "TỔNG SỐ PCS HOÀN THÀNH TRONG THÁNG",
        "ppt_shape": "tong_so_pcs_shape"
    }

]

BLOCK_CONFIG = [
    {
        "name": "ban_ve_cd",
        "title_prefix": "BẢN VẼ",
        "unit": "Bản vẽ",
        "rows": {
            "ca1": "Thực tích Bản vẽ HT CĐ Ca 1 (Ngày)",
            "ca2": "Thực tích Bản vẽ HT CĐ Ca 2 (Ngày)",
            "mt": "Mục tiêu Bản vẽ  HT CĐ (Ngày)",
            "lk1": "Mục tiêu Bản vẽ HT CĐ Ca 1 (Còn lại/Ngày)",
            "lk2": "Mục tiêu Bản vẽ HT CĐ Ca 2(Còn lại/Ngày)",
            "m_ca1": "Thực tích Bản vẽ HT CĐ Ca 1  (Tháng)",
            "m_ca2": "Thực tích Bản vẽ HT CĐ Ca 2  (Tháng)",
            "m_mt": "Mục tiêu Bản vẽ  HT CĐ  (Tháng)"
        },
            "ppt_shape": "ban_ve_cd_shape"
    },
    {
        "name": "so_pcs_cd",
        "title_prefix": "SỐ PCS",
        "unit": "Số PCS",
        "rows": {
            "ca1": "Thực tích số PCS HT CĐ Ca 1 (Ngày)",
            "ca2": "Thực tích số PCS HT CĐ Ca 2 (Ngày)",
            "mt": "Mục tiêu số PCS  HT CĐ (Ngày)",
            "lk1": "Mục tiêu số PCS HT CĐ Ca 1(Còn lại/Ngày)",
            "lk2": "Mục tiêu số PCS HT CĐ Ca 2(Còn lại/Ngày)",
            "m_ca1": "Thực tích số PCS HT CĐ Ca 1  (Tháng)",
            "m_ca2": "Thực tích số PCS HT CĐ Ca 2  (Tháng)",
            "m_mt": "Mục tiêu số PCS  HT CĐ  (Tháng)"
        },
            "ppt_shape": "So_pcs_cd_shape"
    }
]

PPT_ANCHORS = {
    "global": {
        "slide": 1,
        "tong_ban_ve": "anchor_global_left",
        "tong_so_pcs": "anchor_global_right",
    },
    "block": {
        "slide": 3,
        "ban_ve_cd": "anchor_block_left",
        "so_pcs_cd": "anchor_block_right",
    },
    "year": {
        "slide": 0,
        "year_tsx": "anchor_year_bv_pcs"
    },
    "costdown": {
        "slide": 18,
        "TT_noi_tac": "anchor_costdown_noitac",   
        "TT_tong": "anchor_costdown_tong" 
    }
}
BLOCK_SLIDE_MAP = {
    "AA": 4,   # Cắt vật liệu
    "MA": 5,   # Máy Phay
    "LA": 6,   # Máy Tiện
    "BJ": 7,   # Phay chính xác
    "DA": 8,   # Máy Khoan
    "GS": 9,   # Mài Phẳng
    "GR": 10,  # Mài Tròn Có Tâm
    "GC": 11,  # Mài Tròn Vô Tâm
    "GJ": 12,  # Mài Dụng Cụ
    "MC": 13,  # Phay CNC
    "LN": 14,  # Tiện CNC
    "EN": 15,  # Phóng điện
    "EW": 16,  # Cắt dây
    "GP": 17   # Mài mô phỏng
}

# ========== HELPERS ==========
def load_raw_excel():
    return pd.read_excel(EXCEL_FILE, sheet_name=EXCEL_SHEET, header=None)

def split_blocks(df):
    blocks = []
    idxs = df.index[df[0] == "Công đoạn:"].tolist()

    for i in range(len(idxs)):
        start = idxs[i]
        end = idxs[i+1] if i+1 < len(idxs) else len(df)

        block = df.iloc[start:end].reset_index(drop=True)
        blocks.append(block)

    return blocks


def find_row(df, text):
    # exact
    matches = df.index[df[0] == text].tolist()
    if matches:
        return matches[0]
    # substring fuzzy
    for i, val in enumerate(df[0]):
        if isinstance(val, str) and text.lower() in val.lower():
            return i
    # token-based fuzzy
    tokens = [t.strip() for t in text.split() if t.strip()]
    for i, val in enumerate(df[0]):
        if not isinstance(val, str):
            continue
        low = val.lower()
        if len(tokens) >= 1 and all(tok.lower() in low for tok in tokens[:3]):
            return i
    raise ValueError(f"❌ Không tìm thấy: {text}")
def find_row_exact(df, text):
    text_n = normalize_text(text)
    for i, val in enumerate(df[0]):
        if normalize_text(val) == text_n:
            return i
    raise ValueError(f"❌ Không tìm thấy chính xác dòng (sau normalize): {text}")

def normalize_text(s):
    if not isinstance(s, str):
        return s
    return " ".join(
        s.replace("\xa0", " ")  # NBSP → space
         .strip()
         .upper()
         .split()               # 👈 GỘP NHIỀU SPACE
    )



# ========== BUILD DATA ==========
def build_chart_df_global(df_block, rows_dict):
    i_date = find_row_exact(df_block, ROW_DATE)

    dates = pd.to_datetime(
        df_block.loc[i_date, 1:], dayfirst=True, errors="coerce"
    )

    data = {"Ngày": dates}

    for key, row_name in rows_dict.items():
        idx = find_row_exact(df_block, row_name)
        data[key] = pd.to_numeric(
            df_block.loc[idx, 1:], errors="coerce"
        )

    df = (
        pd.DataFrame(data)
        .dropna(subset=["Ngày"])
        .reset_index(drop=True)
    )

    df["X"] = df.index
    df["Label"] = df["Ngày"].dt.strftime("%m-%d")

    df["BAR"]   = df["tt"]
    df["LINE1"] = df["mt"]
    df["LINE2"] = df["lk"]

    return df


def build_chart_df_block(df_block, rows_dict):
    i_date = find_row(df_block, ROW_DATE)

    # ========= LẤY GIÁ TRỊ GỐC TRONG EXCEL =========
    raw_values = df_block.loc[i_date, 1:].values

    raw_labels = []
    month_idx = []
    day_idx = []

    for i, v in enumerate(raw_values):
        # ===== THÁNG: chỉ khi EXCEL GỐC là text kiểu 1月 / Tháng X =====
        if isinstance(v, str) and (
            "月" in v or v.strip().lower().startswith("tháng")
        ):
            month_idx.append(i)
            # chuẩn hoá label tháng
            lbl = v.replace("月", "").replace("Tháng", "").strip()
            raw_labels.append(f"Tháng {lbl}")

        # ===== NGÀY: mọi thứ còn lại coi là ngày =====
        else:
            day_idx.append(i)
            try:
                d = pd.to_datetime(v)
                raw_labels.append(d.strftime("%m-%d"))
            except Exception:
                raw_labels.append(str(v))

    # ========= DATAFRAME THÁNG =========
    df_month = pd.DataFrame()

    if month_idx and {"m_ca1", "m_ca2", "m_mt"}.issubset(rows_dict):
        idx_m_ca1 = find_row(df_block, rows_dict["m_ca1"])
        idx_m_ca2 = find_row(df_block, rows_dict["m_ca2"])
        idx_m_mt  = find_row(df_block, rows_dict["m_mt"])

        df_month = pd.DataFrame({
            "Label": [raw_labels[i] for i in month_idx],
            "Ca1": pd.to_numeric(
                df_block.loc[idx_m_ca1, 1:].values, errors="coerce"
            )[month_idx],
            "Ca2": pd.to_numeric(
                df_block.loc[idx_m_ca2, 1:].values, errors="coerce"
            )[month_idx],
            "Mục tiêu": pd.to_numeric(
                df_block.loc[idx_m_mt, 1:].values, errors="coerce"
            )[month_idx],
        }).fillna(0)

        df_month["X"] = range(len(df_month))

    # ========= DATAFRAME NGÀY =========
    idx_ca1 = find_row(df_block, rows_dict["ca1"])
    idx_ca2 = find_row(df_block, rows_dict["ca2"])
    idx_mt  = find_row(df_block, rows_dict["mt"])
    idx_lk1 = find_row(df_block, rows_dict["lk1"])
    idx_lk2 = find_row(df_block, rows_dict["lk2"])

    df_daily = pd.DataFrame({
        "Label": [raw_labels[i] for i in day_idx],
        "Ca1": pd.to_numeric(
            df_block.loc[idx_ca1, 1:].values, errors="coerce"
        )[day_idx],
        "Ca2": pd.to_numeric(
            df_block.loc[idx_ca2, 1:].values, errors="coerce"
        )[day_idx],
        "Mục tiêu": pd.to_numeric(
            df_block.loc[idx_mt, 1:].values, errors="coerce"
        )[day_idx],
        "Lũy kế Ca1": pd.to_numeric(
            df_block.loc[idx_lk1, 1:].values, errors="coerce"
        )[day_idx],
        "Lũy kế Ca2": pd.to_numeric(
            df_block.loc[idx_lk2, 1:].values, errors="coerce"
        )[day_idx],
    }).fillna(0)

    # daily X = nối sau tháng (Excel-style)
    day_offset = len(df_month)
    df_daily["X"] = range(day_offset, day_offset + len(df_daily))

    if DEBUG:
        print("---- MONTH DF ----")
        print(df_month)
        print("---- DAILY DF ----")
        print(df_daily)

    return df_daily, df_month


# ========== DRAW ==========
def draw_combo_chart_global(df, title, ylabel, output_file, legend_prefix=""):
    def L(text):
        return f"{legend_prefix} - {text}" if legend_prefix else text

    fig, ax = plt.subplots(figsize=(12, 5))

    # ======================
    # BAR
    # ======================
    bars = ax.bar(
        df["X"], df["BAR"],
        width=0.8,
        color="#9c1c7d",
        label=L("Thực tích")
    )

    for bar in bars:
        h = bar.get_height()
        if h > 0:
            ax.text(
                bar.get_x() + bar.get_width() / 2,
                h * 1.01,
                f"{int(h)}",
                ha="center",
                va="bottom",
                fontsize=9
            )

    # ======================
    # LINE
    # ======================
    ax.plot(
        df["X"], df["LINE1"],
        color="green", marker="o",
        linewidth=2,
        label=L("Mục tiêu")
    )

    ax.plot(
        df["X"], df["LINE2"],
        color="red", marker="o",
        linewidth=2,
        label=L("Lũy kế")
    )

    # ======================
    # TITLE
    # ======================
    ax.set_title(title)

    # ======================
    # ✅ Y LABEL ĐƯA LÊN TRÊN
    # ======================
    ax.set_ylabel(ylabel, rotation=0, fontsize=11)
    ax.yaxis.set_label_coords(-0.055, 1.06)
    ax.yaxis.label.set_horizontalalignment("left")

    ax.grid(axis="y", linestyle="--", alpha=0.4)

    step = max(1, len(df) // 10)
    ax.set_xticks(df["X"].iloc[::step])
    ax.set_xticklabels(df["Label"].iloc[::step])

    ax.legend(
        loc="upper center",
        bbox_to_anchor=(0.5, -0.07),
        ncol=3,
        frameon=False
    )

    plt.subplots_adjust(left=0.06, right=0.98, top=0.92, bottom=0.15)
    plt.savefig(output_file, dpi=150)
    plt.close()




def draw_combo_chart_block_excel_style(df_daily, df_month, title, ylabel, output_file, unit=""):
    fig, ax_left = plt.subplots(figsize=(13, 5.5))
    ax_right = ax_left.twinx()

    # Ẩn trục X phụ
    ax_right.xaxis.set_visible(False)

    width = 0.35  # độ rộng cột

    # ================= THÁNG – CLUSTERED COLUMN (TRỤC TRÁI) =================
    if not df_month.empty:
        bars_month_ca1 = ax_left.bar(
            df_month["X"] - width / 2,
            df_month["Ca1"],
            width=width,
            color="#2ca02c",
            label=f"Ca1 (Tháng-{unit})"
        )

        bars_month_ca2 = ax_left.bar(
            df_month["X"] + width / 2,
            df_month["Ca2"],
            width=width,
            color="#1f77b4",
            label=f"Ca2 (Tháng-{unit})"
        )

        # ===== DATA LABEL CHO THỰC TÍCH THÁNG =====
        for bars in [bars_month_ca1, bars_month_ca2]:
            for b in bars:
                h = b.get_height()
                if h > 0:
                    ax_left.text(
                        b.get_x() + b.get_width() / 2,
                        h * 1.01,
                        f"{int(h)}",
                        ha="center",
                        va="bottom",
                        fontsize=8
                    )

        # ===== MỤC TIÊU THÁNG (KHÔNG HIỂN THỊ KHI = 0) =====
        if "Mục tiêu" in df_month.columns:
            target_month = df_month["Mục tiêu"].replace(0, np.nan)
            ax_left.plot(
                df_month["X"],
                target_month,
                color="red",
                marker="o",
                linewidth=2,
                label=f"Mục tiêu (Tháng-{unit})"
            )

        left_max = max(
            df_month[["Ca1", "Ca2", "Mục tiêu"]].max().max(),
            1
        )
        ax_left.set_ylim(0, left_max * LEFT_Y_MULT)

    # ================= NGÀY – CLUSTERED COLUMN (TRỤC PHẢI) =================
    if not df_daily.empty:
        bars_day_ca1 = ax_right.bar(
            df_daily["X"] - width / 2,
            df_daily["Ca1"],
            width=width,
            color="#66c2a5",
            alpha=0.85,
            label=f"Ca1 (Ngày-{unit})"
        )

        bars_day_ca2 = ax_right.bar(
            df_daily["X"] + width / 2,
            df_daily["Ca2"],
            width=width,
            color="#5da5d8",
            alpha=0.85,
            label=f"Ca2 (Ngày-{unit})"
        )

        # ===== DATA LABEL CHO THỰC TÍCH NGÀY =====
        for bars in [bars_day_ca1, bars_day_ca2]:
            for b in bars:
                h = b.get_height()
                if h > 0:
                    ax_right.text(
                        b.get_x() + b.get_width() / 2,
                        h * 1.02,
                        f"{int(h)}",
                        ha="center",
                        va="bottom",
                        fontsize=7
                    )

        # ===== MỤC TIÊU NGÀY (KHÔNG HIỂN THỊ KHI = 0) =====
        if "Mục tiêu" in df_daily.columns:
            target_day = df_daily["Mục tiêu"].replace(0, np.nan)
            ax_right.plot(
                df_daily["X"],
                target_day,
                color="orange",
                linestyle="--",
                linewidth=1.8,
                marker="o",
                label=f"Mục tiêu (Ngày-{unit})"
            )
        # ===== MỤC TIÊU LŨY KẾ CA 1 & CA 2 (NGÀY) =====
        if "Lũy kế Ca1" in df_daily.columns:
            target_lk_ca1 = df_daily["Lũy kế Ca1"].replace(0, np.nan)
            ax_right.plot(
                df_daily["X"],
                target_lk_ca1,
                color="purple",
                linewidth=2.0,
                linestyle="-",
                label=f"Mục tiêu lũy kế Ca1 (Ngày-{unit})"
            )

        if "Lũy kế Ca2" in df_daily.columns:
            target_lk_ca2 = df_daily["Lũy kế Ca2"].replace(0, np.nan)
            ax_right.plot(
                df_daily["X"],
                target_lk_ca2,
                color="brown",
                linewidth=2.0,
                linestyle="-",
                label=f"Mục tiêu lũy kế Ca2 (Ngày-{unit})"
            )


        right_max = max(
            df_daily[["Ca1", "Ca2", "Mục tiêu", "Lũy kế Ca1", "Lũy kế Ca2"]].max().max(),
            1
        )
        ax_right.set_ylim(0, right_max * RIGHT_Y_MULT)

    # ================= CATEGORY X =================
    labels = df_month["Label"].tolist() + df_daily["Label"].tolist()
    positions = list(range(len(labels)))

    STEP = 3  # hiển thị 1 label / 3 ngày
    show_pos = positions[::STEP]
    show_labels = labels[::STEP]

    ax_left.xaxis.set_major_locator(FixedLocator(show_pos))
    ax_left.xaxis.set_major_formatter(FixedFormatter(show_labels))
    ax_left.set_xlim(-0.5, len(labels) - 0.5)

    plt.setp(ax_left.get_xticklabels(), rotation=0, ha="center")

    # ================= LEGEND =================
    h1, l1 = ax_left.get_legend_handles_labels()
    h2, l2 = ax_right.get_legend_handles_labels()


    ax_left.legend(
        h1 + h2,
        l1 + l2,
        loc="lower center",
        bbox_to_anchor=(0.5, -0.19),
        ncol=4,
        frameon=False,
        fontsize=9
    )


    ax_left.grid(axis="y", linestyle="--", alpha=0.3)
    # ======================
    # AXIS LABEL – ĐƯA LÊN TRÊN
    # ======================

    # Trục trái (Tháng)
    ax_left.set_ylabel("Số lượng", rotation=0, fontsize=11)
    ax_left.yaxis.set_label_coords(-0.05, 1.06)
    ax_left.yaxis.label.set_horizontalalignment("left")

    # Trục phải (Ngày)
    ax_right.set_ylabel("Số lượng (Ngày)", rotation=0, fontsize=11)
    ax_right.yaxis.set_label_coords(1.05, 1.06)
    ax_right.yaxis.label.set_horizontalalignment("right")

    plt.title(title)
    plt.subplots_adjust(left=0.06,right=0.95,top=0.92,bottom=0.15)
    plt.savefig(output_file, dpi=150)
    plt.close()

def draw_costdown_chart(df, title, output_file):
    bars, lines, months = classify_costdown_rows(df)

    x = np.arange(len(months))
    fig, ax = plt.subplots(figsize=(13, 5.5))

    # ======================
    # STACKED BAR – NỘI TÁC
    # ======================
    bottom = np.zeros(len(months))
    colors = ["#8dd3c7", "#fb8072", "#80b1d3"]

    for i, vals in enumerate(bars):
        ax.bar(
            x,
            vals,
            bottom=bottom,
            width=0.55,
            color=colors[i],
            label=f"Nội tác {i+1}"
        )
        bottom += vals

    # ======================
    # LINE – LŨY KẾ
    # ======================
    ax2 = ax.twinx()

    line_colors = ["#e41a1c", "#377eb8"]
    for i, vals in enumerate(lines):
        ax2.plot(
            x,
            vals,
            marker="o",
            linewidth=2.5,
            color=line_colors[i],
            label=f"Lũy kế {i+1}"
        )

        # Data label line
        for xi, yi in zip(x, vals):
            if yi > 0:
                ax2.text(
                    xi, yi,
                    f"{int(yi)}",
                    ha="center",
                    va="bottom",
                    fontsize=9,
                    color=line_colors[i]
                )

    # ======================
    # AXIS & TITLE
    # ======================
    ax.set_xticks(x)
    ax.set_xticklabels(months)

    ax.set_ylabel("Giá trị nội tác", rotation=0, fontsize=11)
    ax.yaxis.set_label_coords(-0.05, 1.06)

    ax2.set_ylabel("Lũy kế", rotation=0, fontsize=11)
    ax2.yaxis.set_label_coords(1.05, 1.06)

    fig.suptitle(title, fontsize=16, fontweight="bold")

    ax.grid(axis="y", linestyle="--", alpha=0.4)

    # ======================
    # LEGEND
    # ======================
    h1, l1 = ax.get_legend_handles_labels()
    h2, l2 = ax2.get_legend_handles_labels()

    ax.legend(
        h1 + h2,
        l1 + l2,
        loc="lower center",
        bbox_to_anchor=(0.5, -0.22),
        ncol=5,
        frameon=False
    )

    plt.tight_layout(rect=[0, 0, 1, 0.92])
    plt.savefig(output_file, dpi=150)
    plt.close()


# ========== PPT ==========
def replace_image_in_ppt(prs, shape_name, image_path, slide_index):
    slide = prs.slides[slide_index]

    print(f"🔎 Tìm anchor '{shape_name}' trong slide {slide_index}")

    found = False

    for shape in slide.shapes:
        print("   └─ shape.name =", repr(shape.name))

        if shape.name == shape_name:
            found = True

            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height

            slide.shapes._spTree.remove(shape._element)

            pic = slide.shapes.add_picture(
                str(image_path),
                left,
                top,
                width=width,
                height=height
            )
            pic.name = shape_name

            # prs.save(PPT_FILE)
            print(f"✅ Đã update ảnh tại anchor {shape_name}")
            return

    if not found:
        print(f"❌ KHÔNG tìm thấy anchor '{shape_name}' trong slide {slide_index}")  
         
from pptx.util import Pt
from pptx.dml.color import RGBColor

def update_title_in_ppt(
    prs,
    slide_index,
    title_text,
    title_shape_name="title_block",
    font_name="Calibri",
    font_size_pt=20,
    bold=True,
    color_rgb=(0, 0, 0)
):
    slide = prs.slides[slide_index]

    for shape in slide.shapes:
        if shape.name == title_shape_name and shape.has_text_frame:
            tf = shape.text_frame
            tf.clear()

            p = tf.paragraphs[0]
            p.text = title_text
            p.font.name = font_name
            p.font.size = Pt(font_size_pt)
            p.font.bold = bold
            p.font.italic = False
            p.font.color.rgb = RGBColor(*color_rgb)
            return

    print(f"⚠️ Không tìm thấy title '{title_shape_name}' trong slide {slide_index}")

# =================================================
# YEAR SUMMARY – BV + PCS COMBO (A1:I13)
# =================================================
def load_year_summary():
    df = pd.read_excel(
        EXCEL_YEAR_FILE,
        sheet_name=EXCEL_YEAR_SHEET,
        usecols="A:I",
        nrows=13
    )

    df.columns = [
        "Tháng",
        "Ngày làm việc",
        "MT Giờ",
        "MT BV",
        "TT BV",
        "P/MH BV",
        "MT PCS",
        "TT PCS",
        "P/MH PCS"
    ]

    df = df[df["Tháng"].notna()].reset_index(drop=True)
    # ✅ CHUẨN HOÁ: 1月 → Tháng 1
    df["Tháng"] = df["Tháng"].astype(str).str.replace("月", "", regex=False)
    df["Tháng"] = "Tháng " + df["Tháng"]

    return df


def draw_year_bv_pcs_combo_chart(df, title, output_file):
    x = range(len(df))
    width = 0.35

    fig, ax = plt.subplots(figsize=(13, 5))
    ax2 = ax.twinx()

    # ======================
    # BAR – BẢN VẼ
    # ======================
    bars_mt_bv = ax.bar(
        [i - width/2 for i in x],
        df["MT BV"],
        width=width,
        color="#a6cee3",
        label="Mục tiêu Bản vẽ"
    )

    bars_tt_bv = ax.bar(
        [i + width/2 for i in x],
        df["TT BV"],
        width=width,
        color="#9c1c7d",
        label="Thực tích Bản vẽ"
    )

    # ======================
    # DATA LABEL – BẢN VẼ (CẢ MT & TT ĐỀU GIỮA CỘT)
    # ======================
    for b in bars_mt_bv:
        h = b.get_height()
        if h > 0:
            ax.text(
                b.get_x() + b.get_width() / 2,
                h / 2,
                f"{int(h)}",
                ha="center",
                va="center",
                fontsize=8,
                color="#555555"      # xám trung tính – phân biệt mục tiêu
            )

    for b in bars_tt_bv:
        h = b.get_height()
        if h > 0:
            ax.text(
                b.get_x() + b.get_width() / 2,
                h / 2,
                f"{int(h)}",
                ha="center",
                va="center",
                fontsize=9,
                fontweight="bold",
                color="#000000"      # đen đậm – nhấn mạnh thực tích
            )
                

    # ======================
    # LINE – PCS
    # ======================
    line_mt_pcs = ax2.plot(
        x,
        df["MT PCS"],
        color="orange",
        marker="o",
        linewidth=2,
        label="Mục tiêu PCS"
    )[0]

    line_tt_pcs = ax2.plot(
        x,
        df["TT PCS"],
        color="#1f78b4",
        marker="o",
        linewidth=2,
        label="Thực tích PCS"
    )[0]

    # ✅ DATA LABEL – LINE (CENTER)
    for i, y in enumerate(df["MT PCS"]):
        if y > 0:
            text_transform = offset_copy(
                ax2.transData, fig=fig, x=0, y=6, units='points'
            )
            ax2.text(
                x[i],
                y,
                f"{int(y)}",
                transform=text_transform,
                ha="center",
                va="bottom",
                fontsize=8,
                color="orange"
            )

    for i, y in enumerate(df["TT PCS"]):
        if y > 0:
            text_transform = offset_copy(
                ax2.transData, fig=fig, x=0, y=-8, units='points'
            )
            ax2.text(
                x[i],
                y,
                f"{int(y)}",
                transform=text_transform,
                ha="center",
                va="top",
                fontsize=8,
                color="#1f78b4"
            )

    # ======================
    # TITLE
    # ======================
    fig.suptitle(
        title,
        fontsize=16,
        fontweight="bold",
        y=0.98
    )

    # ======================
    # AXIS & LABEL
    # ======================
    ax.set_xticks(x)
    ax.set_xticklabels(df["Tháng"])

    ax.set_ylabel("Số Bản vẽ", fontsize=11, rotation=0)
    ax.yaxis.set_label_coords(-0.045, 1.08)
    ax.yaxis.label.set_horizontalalignment("left")

    ax2.set_ylabel("Số PCS", fontsize=11, rotation=0)
    ax2.yaxis.set_label_coords(1.045, 1.08)
    ax2.yaxis.label.set_horizontalalignment("right")

    ax.grid(axis="y", linestyle="--", alpha=0.4)

    # ======================
    # LEGEND
    # ======================
    h1, l1 = ax.get_legend_handles_labels()
    h2, l2 = ax2.get_legend_handles_labels()

    ax.legend(
        h1 + h2,
        l1 + l2,
        loc="upper center",
        bbox_to_anchor=(0.5, -0.15),
        ncol=4,
        frameon=False
    )

    plt.tight_layout(rect=[0, 0, 1, 0.94])
    plt.savefig(output_file, dpi=150)
    plt.close()

def process_year_summary_combo(prs):
    df = load_year_summary()

    img_path = OUT_DIR / "year_summary_bv_pcs_combo.png"

    draw_year_bv_pcs_combo_chart(
        df,
        "TỔNG BẢN VẼ & PCS HOÀN THÀNH THEO THÁNG – NĂM 2026",
        img_path
    )

    slide_idx = PPT_ANCHORS["year"]["slide"]

    # ✅ THÊM DÒNG NÀY
    update_title_in_ppt(
        prs,
        slide_idx,
        "TỔNG BẢN VẼ & PCS HOÀN THÀNH THEO THÁNG – NĂM 2026",
        title_shape_name="title_year",
        font_name="Calibri",
        font_size_pt=20,
        bold=True,
        color_rgb=(0, 0, 128)
    )

    replace_image_in_ppt(
        prs,
        PPT_ANCHORS["year"]["year_tsx"],
        img_path,
        slide_index=slide_idx
    )

def load_file3_G1_L16():
    df = pd.read_excel(
        "ThucTich_Data.xlsx",
        sheet_name=0,
        usecols="G:L",
        skiprows=2,      # 🔥 BỎ DÒNG TIÊU ĐỀ
        nrows=13
    )


    df.columns = [
        "Công đoạn",
        "Phụ tải CĐ hiện tại",
        "Phụ tải CĐ tổng",
        "Số bản vẽ trễ hẹn",
        "Số bản vẽ hoàn thành",
        "Hiệu suất (Ngày)"
    ]

    return df

from pptx.util import Inches, Pt

def replace_table_by_anchor_centered(
    prs,
    slide_index,
    anchor_name,
    df,
    report_date
):
    slide = prs.slides[slide_index]

    for shape in slide.shapes:
        if shape.name == anchor_name:

            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            slide.shapes._spTree.remove(shape._element)

            cols = df.shape[1]
            data_rows = len(df)

            table_shape = slide.shapes.add_table(
                data_rows + 2,   # title + header + data
                cols,
                left, top, width, height
            )
            table_shape.name = anchor_name
            table = table_shape.table

            # ===== COLUMN WIDTH =====
            for c in range(cols):
                table.columns[c].width = int(width / cols)

            # ==================================================
            # TITLE
            # ==================================================
            title_cell = table.cell(0, 0)
            title_cell.merge(table.cell(0, cols - 1))
            title_cell.text = (
                f"Bản vẽ Hoàn thành, Trễ hẹn và Phụ tải Công đoạn – {report_date}"
            )
            title_cell.fill.solid()
            title_cell.fill.fore_color.rgb = RGBColor(21, 96, 130)

            for p in title_cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                p.font.size = Pt(16)
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)

            # ==================================================
            # HEADER
            # ==================================================
            for c, col_name in enumerate(df.columns):
                cell = table.cell(1, c)
                cell.text = col_name
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(248, 203, 173)

                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
                    p.font.size = Pt(15)
                    p.font.bold = True

            # ==================================================
            # DATA
            # ==================================================
            for r in range(data_rows):
                for c in range(cols):
                    val = df.iloc[r, c]

                    if isinstance(val, float):
                        txt = f"{val:.2f}"
                    else:
                        txt = str(val)

                    cell = table.cell(r + 2, c)
                    cell.text = txt
                    cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

                    for p in cell.text_frame.paragraphs:
                        p.alignment = PP_ALIGN.CENTER
                        p.font.size = Pt(15)

                        if df.columns[c] == "Số bản vẽ trễ hẹn":
                            p.font.color.rgb = RGBColor(255, 0, 0)

                    cell.fill.solid()
                    cell.fill.fore_color.rgb = (
                        RGBColor(226, 234, 248)
                        if r % 2 == 0
                        else RGBColor(255, 255, 255)
                    )

            print("✅ Table tạo GIỐNG HỆT HÌNH MẪU")
            return

    print(f"❌ Không tìm thấy anchor '{anchor_name}'")
def update_table_title_by_anchor(
    prs,
    slide_index,
    anchor_name,
    title_text,
    font_name="Calibri",
    font_size_pt=14,
    bold=True,
    color_rgb=(0, 0, 128)
):
    slide = prs.slides[slide_index]

    for shape in slide.shapes:
        if shape.name == anchor_name and shape.has_text_frame:
            tf = shape.text_frame
            tf.clear()

            p = tf.paragraphs[0]
            p.text = title_text
            p.font.name = font_name
            p.font.size = Pt(font_size_pt)
            p.font.bold = bold
            p.font.color.rgb = RGBColor(*color_rgb)

            print(f"✅ Đã update tiêu đề table tại '{anchor_name}'")
            return

    print(f"❌ Không tìm thấy anchor tiêu đề '{anchor_name}'") 

    
def load_costdown_B3_J27():
    df = pd.read_excel(
        EXCEL_COSTDOWN_FILE,
        sheet_name=EXCEL_COSTDOWN_SHEET,
        usecols="B:J",   # B → J
        skiprows=2,      # bắt đầu từ B3
        nrows=25         # đến J27
    )

    df = df.dropna(how="all")

    # Cột đầu tiên là tên hạng mục
    df.rename(columns={df.columns[0]: "Hạng mục"}, inplace=True)

    return df
def classify_costdown_rows(df):
    """
    Phân loại các dòng trong B3:J27 dựa trên đặc trưng số học
    Trả về:
        bars: list 3 mảng (dùng stacked bar)
        lines: list 2 mảng (dùng line lũy kế)
    """
    months = [c for c in df.columns if isinstance(c, str) and c.startswith("T")]
    data_rows = []

    for idx, row in df.iterrows():
        vals = pd.to_numeric(row[months], errors="coerce").fillna(0).values
        data_rows.append((idx, vals))

    # === Lũy kế: luôn tăng hoặc giữ nguyên ===
    cumulative_rows = []
    normal_rows = []

    for idx, vals in data_rows:
        if np.all(np.diff(vals) >= 0) and vals.sum() > 0:
            cumulative_rows.append(vals)
        else:
            normal_rows.append(vals)

    if len(cumulative_rows) < 2:
        raise ValueError("❌ Không xác định được đủ 2 dòng lũy kế")

    if len(normal_rows) < 3:
        raise ValueError("❌ Không xác định được đủ 3 dòng nội tác")

    # === Quy ước ===
    bars = normal_rows[:3]       # 3 dòng nội tác
    lines = cumulative_rows[:2]  # 2 dòng lũy kế

    return bars, lines, months

def process_costdown(prs):
    df = load_costdown_B3_J27()

    img_path = OUT_DIR / "costdown_noi_tac.png"

    draw_costdown_chart(
        df,
        "COSTDOWN NỘI TÁC – KẾ HOẠCH & THỰC TẾ NĂM 2026",
        img_path
    )

    slide_idx = PPT_ANCHORS["costdown"]["slide"]

    # Update title
    update_title_in_ppt(
        prs,
        slide_idx,
        "COSTDOWN NỘI TÁC – NĂM 2026",
        title_shape_name="title_costdown",
        font_size_pt=20,
        bold=True,
        color_rgb=(0, 0, 128)
    )

    # Replace chart image
    replace_image_in_ppt(
        prs,
        PPT_ANCHORS["costdown"]["TT_noi_tac"],
        img_path,
        slide_index=slide_idx
    )

# ========== MAIN ==========
def main():
    # ✅ 1. LOAD EXCEL
    df_raw = load_raw_excel()

    # ✅ LOAD PPT 1 LẦN
    prs = Presentation(PPT_FILE)

    # =================================================
    # GLOBAL charts
    # =================================================
    # ✅ LẤY SLIDE GLOBAL
    slide_idx = PPT_ANCHORS["global"]["slide"]
    title_text = "MỤC TIÊU – THỰC TÍCH SẢN LƯỢNG TOÀN BỘ CÔNG ĐOẠN"
    update_title_in_ppt(prs, slide_idx, title_text, title_shape_name="title_global",font_name="Calibri",font_size_pt=20,bold=True,color_rgb=(0, 0, 128))   # xanh đậm
    for cfg in GLOBAL_CONFIG:
        try:
            df_chart = build_chart_df_global(df_raw, cfg["rows"])
        except Exception as e:
            print(f"⚠️ Bỏ qua GLOBAL {cfg['name']} vì lỗi: {e}")
            continue

        img_path = OUT_DIR / f"{cfg['name']}.png"
        draw_combo_chart_global(
            df=df_chart,
            title=cfg.get("title", cfg["name"]),
            ylabel="Số lượng",
            output_file=img_path,
            legend_prefix=cfg.get("legend_prefix", "")
        )
        


        # ✅ LẤY ANCHOR
        anchor = PPT_ANCHORS["global"][cfg["name"]]

        # ✅ GẮN ẢNH VÀO PPT
        replace_image_in_ppt(
            prs,
            anchor,
            img_path,
            slide_index=slide_idx
        )

 

    # =================================================
    # BLOCK charts
    # =================================================
    blocks = split_blocks(df_raw)   # ✅ df_raw đã tồn tại

    def get_cong_doan_name(block):
        """
        Excel:
        A: Công đoạn:
        B: AA / MA / LA ...
        """
        for i, v in enumerate(block[0]):
            if str(v).strip() == "Công đoạn:":
                cd_code = str(block.loc[i, 1]).strip().upper()
                return cd_code
        return None
  
    print("===== DUMP SHAPES PPT =====")
    for idx, slide in enumerate(prs.slides):
        print(f"\nSlide index {idx}:")
        for shape in slide.shapes:
            print("   ", repr(shape.name))

    for i, block in enumerate(blocks):
        cd_code = get_cong_doan_name(block)

        if not cd_code:
            print("⏭️ Bỏ qua block không có Công đoạn")
            continue

        print(f"📌 Công đoạn {cd_code}")

        if cd_code not in BLOCK_SLIDE_MAP:
            print(f"⚠️ Không có slide cho công đoạn {cd_code}")
            continue

        slide_idx = BLOCK_SLIDE_MAP[cd_code]
        title_text = f"TÍNH SẢN XUẤT – TỈ LỆ HOẠT ĐỘNG MÁY CÔNG ĐOẠN {cd_code}"
        
        update_title_in_ppt(
            prs,
            slide_idx,
            title_text,
            title_shape_name="title_block",
            font_name="Calibri",
            font_size_pt=20,
            bold=True,
            color_rgb=(0, 0, 128)   # xanh đậm
        )


        for cfg in BLOCK_CONFIG:
            df_daily, df_month = build_chart_df_block(block, cfg["rows"])

            img_path = OUT_DIR / f"{cfg['name']}_{i+1}.png"

            draw_combo_chart_block_excel_style(
                df_daily=df_daily,
                df_month=df_month,
                title=f"{cfg['title_prefix']} – CĐ {cd_code} – CA 1 + CA 2",
                ylabel="Số lượng",
                output_file=img_path,
                unit=cfg["unit"]
            )

            anchor = PPT_ANCHORS["block"][cfg["name"]]

            replace_image_in_ppt(
                prs,
                anchor,
                img_path,
                slide_index=slide_idx
            )

    process_year_summary_combo(prs)
    # =================================================
    # FILE 3 – TABLE (G1:L16) → SLIDE 3
    # =================================================
    df_file3 = load_file3_G1_L16()

    # ✅ TIÊU ĐỀ TABLE (NẰM TRÊN)
    update_table_title_by_anchor(
        prs,
        slide_index=2,
        anchor_name="title_file3_table",
        title_text="PHỤ TẢI CÔNG ĐOẠN – THỰC TÍCH",
        font_size_pt=20
    )


    report_date = datetime.now().strftime("%Y-%m-%d")

    replace_table_by_anchor_centered(
        prs,
        slide_index=2,
        anchor_name="anchor_file3_table",
        df=df_file3,
        report_date=report_date
    )

    process_costdown(prs)

    prs.save(PPT_FILE)

    print("🎉 HOÀN TẤT")

if __name__ == "__main__":
    main()
