import pandas as pd
import matplotlib.pyplot as plt
import numpy as np
# import matplotlib.dates as mdates
from pptx import Presentation
from pathlib import Path
from matplotlib.ticker import FixedLocator, FixedFormatter
from matplotlib.transforms import offset_copy
plt.rcParams["font.family"] = "Segoe UI"
plt.rcParams["axes.unicode_minus"] = False
# hoặc:
# plt.rcParams["font.family"] = "Yu Gothic"
# plt.rcParams["font.family"] = "Noto Sans CJK JP"
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from datetime import datetime
from matplotlib.patches import Patch
from matplotlib.lines import Line2D
from pptx.util import Pt
from pptx.dml.color import RGBColor
from datetime import datetime, timedelta
import tempfile
import shutil
import os
import stat
import subprocess
import time

def copy_to_server_safe(src: Path, dst_dir: Path, retry=3):
    """
    Copy file PNG lên UNC server an toàn:
    - Không dùng copy2 (tránh metadata + permission)
    - Dùng file tạm + os.replace (atomic)
    - Retry để né file lock ngắn hạn
    """
    dst_final = dst_dir / src.name
    dst_tmp = dst_dir / (src.stem + "__tmp" + src.suffix)

    for _ in range(retry):
        try:
            shutil.copyfile(src, dst_tmp)   # ✅ nội dung thuần
            os.replace(dst_tmp, dst_final)  # ✅ atomic overwrite
            return
        except PermissionError:
            time.sleep(1)

    raise PermissionError(f"❌ Không thể ghi file lên server: {dst_final}")


# =========================
# GLOBAL CHART THEME
# =========================
CHART_THEME = {
    # ----- FONT -----
    "font_family": "Segoe UI",

    # ----- TITLE -----
    "title": {
        "fontsize": 18,
        "fontweight": "bold",
        "y": 0.975,
    },

    # ----- AXIS LABEL -----
    "ylabel": {
        "fontsize": 9,
        "rotation": 0
    },

    # ----- GRID -----
    "grid": {
        "linestyle": "--",
        "alpha": 0.25
    },

    # ----- BAR COLORS -----
    "bar": {
        "month_ca1": "#2ca02c",
        "month_ca2": "#1f77b4",
        "day_ca1": "#66c2a5",
        "day_ca2": "#5da5d8"
    },

    # ----- LINE COLORS -----
    "line": {
        "target_month": "red",
        "target_day": "orange",
        "remain_ca1": "purple",
        "remain_ca2": "brown"
    },

    # ----- LINE STYLES -----
    "linestyle": {
        "target": "-",
        "remain": "--"
    },

    # ----- MARKERS -----
    "marker": {
        "target": "o",
        "remain": "s"
    },

    # ----- LEGEND -----
    "legend": {
        "fontsize": 9,
        "ncol": 4,
        "columnspacing": 1.6,
        "handletextpad": 0.6
    },

    # ----- LAYOUT -----
    "layout": {
        "left": 0.035,
        "right": 0.965,
        "top": 0.88,
        "bottom": 0.15
    }
}

def save_figure(fig, output_file):
    """
    Chuẩn hóa export PNG cho toàn bộ dashboard:
    - Không bbox_inches="tight"
    - Không dính trần title
    - Chuẩn PPT / TV
    """
    fig.savefig(
        output_file.with_suffix(".png"),
        dpi=240,
        pad_inches=0.04
    )
    plt.close(fig)

#========== Chính ==========
EXCEL_FILE = Path(r"\\vdm-fsvr\Du lieu dung chung VDM\PROJECT-IT HOA NGHIEP VU\Báo Cáo Hằng Ngày CK\Hằng ngày_data_view.xlsx")
EXCEL_SHEET = "Hằng ngày"
EXCEL_YEAR_FILE = Path(r"\\vdm-fsvr\Du lieu dung chung VDM\PROJECT-IT HOA NGHIEP VU\Báo Cáo Hằng Ngày CK\Tính sản xuất năm 2026.xlsx")
EXCEL_YEAR_SHEET = "NĂM 2026"
EXCEL_THUCTICH_FILE = Path(r"\\vdm-fsvr\Du lieu dung chung VDM\PROJECT-IT HOA NGHIEP VU\Báo Cáo Hằng Ngày CK\ThucTich_Data.xlsx")
EXCEL_THUCTICH_SHEET = "Sheet1"
EXCEL_COSTDOWN_FILE = Path(r"\\vdm-fsvr\Cokhi-機工\CONG VIEC CHUNG 2026\★2026年度年次目標設定\5.COSTDOWN 2026\5.MỤC TIÊU 2026\COSTDOWN 2026\Costdown nội tác 2026.xlsx")
EXCEL_COSTDOWN_SHEET = "Quản lý MT năm"
EXCEL_TIMER_FILE = Path(r"\\vdm-fsvr\Du lieu dung chung VDM\PROJECT-IT HOA NGHIEP VU\Báo Cáo Hằng Ngày CK\DulieutudongTimer_view.xlsx")
EXCEL_TIMER_SHEET = "Tong_Hop"
PPT_FILE = Path(r"\\vdm-fsvr\Du lieu dung chung VDM\PROJECT-IT HOA NGHIEP VU\Báo Cáo Hằng Ngày CK\Bao_Cao_Hang_Ngay.pptx")
ROW_DATE = "Ngày làm việc"

# ========== Test ==========
# EXCEL_FILE = Path(r"D:\Code_cokhi\Bao_Cao_Tu_Dong_PPTX_New\Excel\Hằng ngày_data.xlsx")
# EXCEL_SHEET = "Hằng ngày"
# EXCEL_YEAR_FILE = Path(r"D:\Code_cokhi\Bao_Cao_Tu_Dong_PPTX_New\Excel\Tính sản xuất năm 2026.xlsx")
# EXCEL_YEAR_SHEET = "NĂM 2026"
# EXCEL_THUCTICH_FILE = Path(r"D:\Code_cokhi\Bao_Cao_Tu_Dong_PPTX_New\Excel\ThucTich_Data.xlsx")
# EXCEL_THUCTICH_SHEET = "Sheet1"
# EXCEL_COSTDOWN_FILE = Path(r"D:\Code_cokhi\Bao_Cao_Tu_Dong_PPTX_New\Excel\Costdown nội tác 2026.xlsx")
# EXCEL_COSTDOWN_SHEET = "Quản lý MT năm"
# EXCEL_TIMER_FILE = Path(r"D:\Code_cokhi\Bao_Cao_Tu_Dong_PPTX_New\Excel\DulieutudongTimer.xlsx")
# EXCEL_TIMER_SHEET = "Tong_Hop"
# PPT_FILE = Path(r"D:\Code_cokhi\Bao_Cao_Tu_Dong_PPTX_New\Bao_Cao_Hang_Ngay.pptx")
# ROW_DATE = "Ngày làm việc"


LOCAL_OUT_DIR = Path(tempfile.gettempdir()) / "BaoCaoHangNgay"
LOCAL_OUT_DIR.mkdir(parents=True, exist_ok=True)

SERVER_OUT_DIR = Path(r"charts")
SERVER_OUT_DIR.mkdir(parents=True, exist_ok=True)

REQUIRED_FILES = {
    "Daily": EXCEL_FILE,
    "Year": EXCEL_YEAR_FILE,
    "ThucTich": EXCEL_THUCTICH_FILE,
    "Costdown": EXCEL_COSTDOWN_FILE,
    "Timer": EXCEL_TIMER_FILE,
    "PPT": PPT_FILE,
}

for name, p in REQUIRED_FILES.items():
    if not p.exists():
        raise FileNotFoundError(f"❌ [{name}] Không tìm thấy file: {p}")
# Scale multipliers (tinh chỉnh để biểu đồ trông giống thủ công)
LEFT_Y_MULT = 1.3   # scale cho trục trái (tháng)
RIGHT_Y_MULT = 1.3  # scale cho trục phải (ngày)
MONTH_WIDTH_DAYS = 18
DAY_BAR_WIDTH = 0.6

# Debug: bật để in df_daily / df_month (False khi chạy chính thức)
DEBUG = False

# ========== CONFIGS ==========
GLOBAL_CONFIG = [
    {
        "name": "tong_ban_ve",
        "legend_prefix": "Bản vẽ",
        "rows": {
            "tt": "THỰC TÍCH SỐ BẢN VẼ HOÀN THÀNH TRONG NGÀY",
            "mt": "MỤC TIÊU SỐ BẢN VẼ HOÀN THÀNH TRONG NGÀY",
            "lk": "MỤC TIÊU LŨY KẾ BẢN VẼ HT TRONG NGÀY ( CÒN LẠI)",
        },
        "title": "TỔNG BẢN VẼ HOÀN THÀNH TRONG NGÀY",
        "ppt_shape": "tong_ban_ve_shape"
    },
    
    {
        "name": "tong_so_pcs",
        "legend_prefix": "Số PCS",
        "rows": {
            "tt": "THỰC TÍCH SỐ PCS HOÀN THÀNH TRONG NGÀY",
            "mt": "MỤC TIÊU SỐ PCS HOÀN THÀNH TRONG NGÀY",
            "lk": "MỤC TIÊU LŨY KẾ SỐ PCS HT TRONG NGÀY ( CÒN LẠI)",
        },
        "title": "TỔNG SỐ PCS HOÀN THÀNH TRONG NGÀY",
        "ppt_shape": "tong_so_pcs_shape"
    }

]

BLOCK_CONFIG = [
    {
        "name": "ban_ve_cd",
        "title_prefix": "BẢN VẼ",
        "unit": "Bản vẽ",
        "rows": {
            "ca1": "Thực tích Bản vẽ HT CĐ Ca 1 (Ngày)",
            "ca2": "Thực tích Bản vẽ HT CĐ Ca 2 (Ngày)",
            "mt": "Mục tiêu Bản vẽ  HT CĐ (Ngày)",
            "lk1": "Mục tiêu Bản vẽ HT CĐ Ca 1 (Còn lại/Ngày)",
            "lk2": "Mục tiêu Bản vẽ HT CĐ Ca 2(Còn lại/Ngày)",
            "m_ca1": "Thực tích Bản vẽ HT CĐ Ca 1  (Tháng)",
            "m_ca2": "Thực tích Bản vẽ HT CĐ Ca 2  (Tháng)",
            "m_mt": "Mục tiêu Bản vẽ  HT CĐ  (Tháng)"
        },
            "ppt_shape": "ban_ve_cd_shape"
    },
    {
        "name": "so_pcs_cd",
        "title_prefix": "SỐ PCS",
        "unit": "Số PCS",
        "rows": {
            "ca1": "Thực tích số PCS HT CĐ Ca 1 (Ngày)",
            "ca2": "Thực tích số PCS HT CĐ Ca 2 (Ngày)",
            "mt": "Mục tiêu số PCS  HT CĐ (Ngày)",
            "lk1": "Mục tiêu số PCS HT CĐ Ca 1(Còn lại/Ngày)",
            "lk2": "Mục tiêu số PCS HT CĐ Ca 2(Còn lại/Ngày)",
            "m_ca1": "Thực tích số PCS HT CĐ Ca 1  (Tháng)",
            "m_ca2": "Thực tích số PCS HT CĐ Ca 2  (Tháng)",
            "m_mt": "Mục tiêu số PCS  HT CĐ  (Tháng)"
        },
            "ppt_shape": "So_pcs_cd_shape"
    }
]

PPT_ANCHORS = {
    "year": {
        "slide": 0,
        "year_tsx": "anchor_year_bv_pcs"
    },
    "global": {
        "slide": 1,
        "tong_ban_ve": "anchor_global_left",
        "tong_so_pcs": "anchor_global_right",
    },
    "block": {
        "slide": 4,
        "ban_ve_cd": "anchor_block_left",
        "so_pcs_cd": "anchor_block_right",
    },
    "timer": {
        "slide": 3,
        "timer_cd": "anchor_timer_utilization"
    },
    "costdown": {
        "slide": 18,
        "TT_noi_tac": "anchor_costdown_noitac",   
        "TT_tong": "anchor_costdown_tong" 
    }
}
BLOCK_SLIDE_MAP = {
    "AA": 4,   # Cắt vật liệu
    "MA": 5,   # Máy Phay
    "LA": 6,   # Máy Tiện
    "BJ": 7,   # Phay chính xác
    "DA": 8,   # Máy Khoan
    "GS": 9,   # Mài Phẳng
    "GR": 10,  # Mài Tròn Có Tâm
    "GC": 11,  # Mài Tròn Vô Tâm
    "GJ": 12,  # Mài Dụng Cụ
    "MC": 13,  # Phay CNC
    "LN": 14,  # Tiện CNC
    "EN": 15,  # Phóng điện
    "EW": 16,  # Cắt dây
    "GP": 17   # Mài mô phỏng
}

TIMER_SHEETS = [
    "LA", "MA", "BJ", "GC", "GS", "GR",
    "EN", "GJ", "MC", "LN", "EW", "GP"
]
# ========== HELPERS ==========
def load_raw_excel():
    return pd.read_excel(EXCEL_FILE, sheet_name=EXCEL_SHEET, header=None)

def split_blocks(df):
    blocks = []
    idxs = df.index[df[0] == "Công đoạn:"].tolist()

    for i in range(len(idxs)):
        start = idxs[i]
        end = idxs[i+1] if i+1 < len(idxs) else len(df)

        block = df.iloc[start:end].reset_index(drop=True)
        blocks.append(block)

    return blocks


def find_row(df, text):
    # exact
    matches = df.index[df[0] == text].tolist()
    if matches:
        return matches[0]
    # substring fuzzy
    for i, val in enumerate(df[0]):
        if isinstance(val, str) and text.lower() in val.lower():
            return i
    # token-based fuzzy
    tokens = [t.strip() for t in text.split() if t.strip()]
    for i, val in enumerate(df[0]):
        if not isinstance(val, str):
            continue
        low = val.lower()
        if len(tokens) >= 1 and all(tok.lower() in low for tok in tokens[:3]):
            return i
    raise ValueError(f"❌ Không tìm thấy: {text}")
def find_row_exact(df, text):
    text_n = normalize_text(text)
    for i, val in enumerate(df[0]):
        if normalize_text(val) == text_n:
            return i
    raise ValueError(f"❌ Không tìm thấy chính xác dòng (sau normalize): {text}")

def normalize_text(s):
    if not isinstance(s, str):
        return s
    return " ".join(
        s.replace("\xa0", " ")  # NBSP → space
         .strip()
         .upper()
         .split()               # 👈 GỘP NHIỀU SPACE
    )



# ========== BUILD DATA ==========
def build_chart_df_global(df_block, rows_dict):
    i_date = find_row_exact(df_block, ROW_DATE)

    dates = pd.to_datetime(
        df_block.loc[i_date, 1:], dayfirst=True, errors="coerce"
    )

    data = {"Ngày": dates}

    for key, row_name in rows_dict.items():
        idx = find_row_exact(df_block, row_name)
        data[key] = pd.to_numeric(
            df_block.loc[idx, 1:], errors="coerce"
        )

    df = (
        pd.DataFrame(data)
        .dropna(subset=["Ngày"])
        .reset_index(drop=True)
    )

    df["X"] = df.index
    df["Label"] = df["Ngày"].dt.strftime("%m-%d")

    df["BAR"]   = df["tt"]
    df["LINE1"] = df["mt"]
    df["LINE2"] = df["lk"]

    return df


def build_chart_df_block(df_block, rows_dict):
    i_date = find_row(df_block, ROW_DATE)

    # ========= LẤY GIÁ TRỊ GỐC TRONG EXCEL =========
    raw_values = df_block.loc[i_date, 1:].values

    raw_labels = []
    month_idx = []
    day_idx = []

    for i, v in enumerate(raw_values):
        # ===== THÁNG: chỉ khi EXCEL GỐC là text kiểu 1月 / Tháng X =====
        if isinstance(v, str) and (
            "月" in v or v.strip().lower().startswith("tháng")
        ):
            month_idx.append(i)
            # chuẩn hoá label tháng
            lbl = v.replace("月", "").replace("Tháng", "").strip()
            raw_labels.append(f"T{lbl}")

        # ===== NGÀY: mọi thứ còn lại coi là ngày =====
        else:
            day_idx.append(i)
            try:
                d = pd.to_datetime(v)
                raw_labels.append(d.strftime("%m-%d"))
            except Exception:
                raw_labels.append(str(v))

    # ========= DATAFRAME THÁNG =========
    df_month = pd.DataFrame()

    if month_idx and {"m_ca1", "m_ca2", "m_mt"}.issubset(rows_dict):
        idx_m_ca1 = find_row(df_block, rows_dict["m_ca1"])
        idx_m_ca2 = find_row(df_block, rows_dict["m_ca2"])
        idx_m_mt  = find_row(df_block, rows_dict["m_mt"])

        df_month = pd.DataFrame({
            "Label": [raw_labels[i] for i in month_idx],
            "Ca1": pd.to_numeric(
                df_block.loc[idx_m_ca1, 1:].values, errors="coerce"
            )[month_idx],
            "Ca2": pd.to_numeric(
                df_block.loc[idx_m_ca2, 1:].values, errors="coerce"
            )[month_idx],
            "Mục tiêu": pd.to_numeric(
                df_block.loc[idx_m_mt, 1:].values, errors="coerce"
            )[month_idx],
        }).fillna(0)

        df_month["X"] = range(len(df_month))

    # ========= DATAFRAME NGÀY =========
    idx_ca1 = find_row(df_block, rows_dict["ca1"])
    idx_ca2 = find_row(df_block, rows_dict["ca2"])
    idx_mt  = find_row(df_block, rows_dict["mt"])
    idx_lk1 = find_row(df_block, rows_dict["lk1"])
    idx_lk2 = find_row(df_block, rows_dict["lk2"])

    df_daily = pd.DataFrame({
        "Label": [raw_labels[i] for i in day_idx],
        "Ca1": pd.to_numeric(
            df_block.loc[idx_ca1, 1:].values, errors="coerce"
        )[day_idx],
        "Ca2": pd.to_numeric(
            df_block.loc[idx_ca2, 1:].values, errors="coerce"
        )[day_idx],
        "Mục tiêu": pd.to_numeric(
            df_block.loc[idx_mt, 1:].values, errors="coerce"
        )[day_idx],
        "Lũy kế Ca1": pd.to_numeric(
            df_block.loc[idx_lk1, 1:].values, errors="coerce"
        )[day_idx],
        "Lũy kế Ca2": pd.to_numeric(
            df_block.loc[idx_lk2, 1:].values, errors="coerce"
        )[day_idx],
    }).fillna(0)

    # daily X = nối sau tháng (Excel-style)
    day_offset = len(df_month)
    df_daily["X"] = range(day_offset, day_offset + len(df_daily))

    if DEBUG:
        print("---- MONTH DF ----")
        print(df_month)
        print("---- DAILY DF ----")
        print(df_daily)

    return df_daily, df_month


# ========== DRAW ==========
def draw_combo_chart_global(df, title, ylabel, output_file, legend_prefix=""):
    def L(text):
        return f"{legend_prefix} - {text}" if legend_prefix else text


    # ======================
    # TÍNH KPI TỔNG
    # ======================
    total_tt = df["BAR"].sum()
    total_mt = df["LINE1"].sum()
    percent = (total_tt / total_mt * 100) if total_mt > 0 else 0

    # kpi_text = (
    #     f"{legend_prefix}: "
    #     f"{int(total_tt):,}/{int(total_mt):,} "
    #     f"({percent:.1f}%)"
    # )

    # # đổi màu KPI theo % đạt
    # if percent >= 80:
    #     kpi_color = "#2e7d32"   # xanh
    # elif percent >= 60:
    #     kpi_color = "#f9a825"   # vàng
    # else:
    #     kpi_color = "#c62828"   # đỏ

 
    # ======================
    # FIGURE & AXIS
    # ======================
    fig, ax = plt.subplots(figsize=(13,3.5))  # 13 x 7.5 inch

    # ======================
    # BAR – THỰC TÍCH
    # ======================
    bars = ax.bar(
        df["X"], df["BAR"],
        width=0.8,
        color="#209c1c",
        label=f"{legend_prefix} hoàn thành" if legend_prefix else "Hoàn thành"
    )

    # ======================
    # DATA LABEL – THỰC TÍCH (CENTER + WHITE HALO)
    # ======================
    BAR_TEXT_BOX = dict(
        boxstyle="round,pad=0.18",
        facecolor="white",
        alpha=0.75,
        edgecolor="none"
    )

    for bar in bars:
        h = bar.get_height()
        if h > 0:
            ax.text(
                bar.get_x() + bar.get_width() / 2,
                h / 2,                       # ✅ center theo chiều cao cột
                f"{int(h)}",
                ha="center",
                va="center",
                fontsize=9,
                fontweight="normal",
                color="black",
                bbox=BAR_TEXT_BOX
            )

    # ======================
    # LINE – MỤC TIÊU
    # ======================
    ax.plot(
        df["X"], df["LINE1"],
        color=CHART_THEME["line"]["target_month"],
        linestyle=CHART_THEME["linestyle"]["target"],
        marker=CHART_THEME["marker"]["target"],
        linewidth=2,
        label=f"Mục tiêu {legend_prefix} hoàn thành" if legend_prefix else "Mục tiêu hoàn thành"
    )

    # ======================
    # LINE – CÒN LẠI (LŨY KẾ)
    # ======================
    ax.plot(
        df["X"], df["LINE2"],
        color="red",
        linestyle=CHART_THEME["linestyle"]["remain"],
        marker=CHART_THEME["marker"]["remain"],
        linewidth=1.6,
        label=f"Mục tiêu {legend_prefix} còn lại" if legend_prefix else "Mục tiêu còn lại"
    )
    # ======================
    # NÂNG TRẦN TRỤC Y (HEADROOM) – SAFE
    # ======================
    series = pd.concat([
        df["BAR"],
        df["LINE1"],
        df["LINE2"]
    ], axis=0)

    series = series.replace([np.inf, -np.inf], np.nan)
    y_max_data = series.dropna().max()

    if pd.isna(y_max_data) or y_max_data <= 0:
        y_max_data = 1   # ✅ fallback an toàn

    ax.set_ylim(0, y_max_data * 1.18)

    # ======================
    # DATA LABEL – TINH CHỈNH CAO CẤP (RIÊNG CHART NÀY)
    # ======================
    ymax = ax.get_ylim()[1]

    MIN_GAP = ymax * 0.06
    bar_vals = df["BAR"].values

    TEXT_BOX = dict(
        boxstyle="round,pad=0.2",
        facecolor="white",
        alpha=0.85,
        edgecolor="none"
    )

    for x_val, y1, y2 in zip(df["X"], df["LINE1"], df["LINE2"]):
        if y1 <= 0 and y2 <= 0:
            continue

        delta = abs(y1 - y2)
        sep = max(MIN_GAP, delta * 0.25)

        bar_h = bar_vals[int(x_val)]
        extra_bar_gap = ymax * 0.02 if bar_h > y1 * 0.9 else 0
        y_label1 = max(y1 - sep - extra_bar_gap, ymax * 0.06)
        y_label2 = min(y2 + sep, ymax * 0.965)

        # LINE1 – Mục tiêu hoàn thành
        if y1 > 0:
            ax.text(
                x_val,
                y_label1,
                int(y1),
                ha="center",
                va="top",
                fontsize=8,
                color=CHART_THEME["line"]["target_month"],
                bbox=TEXT_BOX
            )

        # LINE2 – Mục tiêu còn lại
        if y2 > 0:
            ax.text(
                x_val,
                y_label2,
                int(y2),
                ha="center",
                va=("bottom" if y_label2 < ymax * 0.96 else "top"),
                fontsize=8,
                color="red",
                bbox=TEXT_BOX
            )
    # ======================
    # X TICK
    # ======================
    step = max(1, len(df) // 10)
    ax.set_xticks(df["X"].iloc[::step])
    ax.set_xticklabels(df["Label"].iloc[::step])

    # ======================
    # GRID
    # ======================
    ax.grid(
        axis="y",
        linestyle=CHART_THEME["grid"]["linestyle"],
        alpha=CHART_THEME["grid"]["alpha"]
    )
    ax.set_axisbelow(True)

    # ======================
    # Y LABEL
    # ======================
    ax.set_ylabel(ylabel, **CHART_THEME["ylabel"])
    ax.yaxis.set_label_coords(-0.025, 1.03)
    ax.yaxis.label.set_horizontalalignment("left")

    # ======================
    # LEGEND
    # ======================
    ax.legend(
        loc="lower center",
        bbox_to_anchor=(0.5, -0.22),
        ncol=3,
        frameon=False,
        fontsize=CHART_THEME["legend"]["fontsize"]
    )

    # ======================
    # TITLE (ĐƯA LÊN CAO HƠN)
    # ======================
    fig.suptitle(
        title,
        fontsize=CHART_THEME["title"]["fontsize"],
        # fontweight=CHART_THEME["title"]["fontweight"],
        y=CHART_THEME["title"]["y"]   # ✅ đã nâng cao
    )
    
    
    y_kpi = 0.94

    # # Nhãn
    # fig.text(
    #     0.70, y_kpi,
    #     f"{legend_prefix} |",
    #     ha="right",
    #     va="center",
    #     fontsize=12,
    #     fontweight="bold",
    #     color="#0c2d57"
    # )

    # Thực tích – xanh đậm
    fig.text(
        0.712, y_kpi,
        f"T/tích: {int(total_tt):,}",
        ha="left",
        va="center",
        fontsize=12,
        fontweight="bold",
        color="#1f77b4"
    )

    # Mục tiêu – đỏ
    fig.text(
        0.798, y_kpi,
        f"M/tiêu: {int(total_mt):,}",
        ha="left",
        va="center",
        fontsize=12,
        fontweight="bold",
        color="#c62828"
    )

    # % hoàn thành – trung tính
    fig.text(
        0.985, y_kpi,
        f"H/thành: {percent:.1f}%",
        ha="right",
        va="center",
        fontsize=12,
        fontweight="bold",
        color="#444444"
    )




    # ======================
    # LAYOUT & SAVE
    # ======================
    plt.subplots_adjust(**CHART_THEME["layout"])
    save_figure(fig, output_file)


def draw_combo_chart_block_excel_style(df_daily, df_month, title, ylabel, output_file, unit=""):
    # ======================
    # FIGURE & AXES
    # ======================
    fig, ax_left = plt.subplots(figsize=(13, 3))
    ax_right = ax_left.twinx()
    ax_right.xaxis.set_visible(False)

    width = 0.45

    # ======================
    # THÁNG – TRỤC TRÁI
    # ======================
    h_ca1_month = h_ca2_month = h_mt_month = None

    if not df_month.empty:
        bars_month_ca1 = ax_left.bar(
            df_month["X"] - width / 2,
            df_month["Ca1"],
            width=width,
            color=CHART_THEME["bar"]["month_ca1"],
            alpha=0.9
        )
        bars_month_ca2 = ax_left.bar(
            df_month["X"] + width / 2,
            df_month["Ca2"],
            width=width,
            color=CHART_THEME["bar"]["month_ca2"],
            alpha=0.9
        )

        h_ca1_month = bars_month_ca1[0]
        h_ca2_month = bars_month_ca2[0]

        # ===== DATA LABEL – HOÀN THÀNH (THÁNG) → CENTER =====
        for bars in (bars_month_ca1, bars_month_ca2):
            for b in bars:
                h = b.get_height()
                if h > 0:
                    ax_left.text(
                        b.get_x() + b.get_width() / 2,
                        h / 2,                     # ✅ CENTER
                        int(h),
                        ha="center",
                        va="center",
                        fontsize=6.5,              # ✅ nhỏ hơn
                        color="black",
                        bbox=dict(
                            boxstyle="round,pad=0.12",
                            facecolor="white",
                            alpha=0.65,
                            edgecolor="none"
                        ),
                        zorder=6
                    )

        if "Mục tiêu" in df_month.columns:
            h_mt_month, = ax_left.plot(
                df_month["X"],
                df_month["Mục tiêu"].replace(0, np.nan),
                color=CHART_THEME["line"]["target_month"],
                linestyle=CHART_THEME["linestyle"]["target"],
                marker=CHART_THEME["marker"]["target"],
                linewidth=2.0
            )

            # ===== DATA LABEL – MỤC TIÊU (THÁNG) =====
            for x_val, y_val in zip(df_month["X"], df_month["Mục tiêu"]):
                if y_val > 0:
                    ax_left.text(
                        x_val,
                        y_val * 1.05,          # ⬆️ đẩy nhẹ lên trên marker
                        int(y_val),
                        ha="center",
                        va="bottom",
                        fontsize=8,
                        color=CHART_THEME["line"]["target_month"]
                    )

        vals = pd.concat([
            df_month["Ca1"],
            df_month["Ca2"],
            df_month["Mục tiêu"]
        ], axis=0).replace([np.inf, -np.inf], np.nan)

        left_max = vals.dropna().max()

        if pd.isna(left_max) or left_max <= 0:
            left_max = 1

        ax_left.set_ylim(0, left_max * LEFT_Y_MULT)

    # ======================
    # NGÀY – TRỤC PHẢI
    # ======================
    h_ca1_day = h_ca2_day = h_mt_day = h_lk_ca1 = h_lk_ca2 = None

    if not df_daily.empty:
        width_day = width * 0.9

        bars_day_ca1 = ax_right.bar(
            df_daily["X"] - width_day / 2,
            df_daily["Ca1"],
            width=width_day,
            color=CHART_THEME["bar"]["day_ca1"],
            alpha=0.85
        )
        bars_day_ca2 = ax_right.bar(
            df_daily["X"] + width_day / 2,
            df_daily["Ca2"],
            width=width_day,
            color=CHART_THEME["bar"]["day_ca2"],
            alpha=0.85
        )

        h_ca1_day = bars_day_ca1[0]
        h_ca2_day = bars_day_ca2[0]

        # ===== DATA LABEL – HOÀN THÀNH (NGÀY) – NÉ VA CA1 / CA2 =====
        ymax_r = ax_right.get_ylim()[1]
        OFFSET = ymax_r * 0.04          # độ lệch nhỏ
        DELTA_MIN = ymax_r * 0.06       # ngưỡng coi là "gần nhau"

        for (bar1, v1), (bar2, v2) in zip(
            zip(bars_day_ca1, df_daily["Ca1"]),
            zip(bars_day_ca2, df_daily["Ca2"])
        ):
            if v1 <= 0 and v2 <= 0:
                continue

            # khoảng cách 2 cột
            delta = abs(v1 - v2)

            # ===== Ca 1 =====
            if v1 > 0:
                if delta < DELTA_MIN:
                    y1 = v1 / 2 - OFFSET   # lệch xuống
                else:
                    y1 = v1 / 2            # center

                ax_right.text(
                    bar1.get_x() + bar1.get_width() / 2,
                    y1,
                    int(v1),
                    ha="center",
                    va="center",
                    fontsize=7,
                    bbox=dict(
                        boxstyle="round,pad=0.12",
                        facecolor="white",
                        alpha=0.65,
                        edgecolor="none"
                    ),
                    zorder=6
                )

            # ===== Ca 2 =====
            if v2 > 0:
                if delta < DELTA_MIN:
                    y2 = v2 / 2 + OFFSET   # lệch lên
                else:
                    y2 = v2 / 2            # center

                ax_right.text(
                    bar2.get_x() + bar2.get_width() / 2,
                    y2,
                    int(v2),
                    ha="center",
                    va="center",
                    fontsize=7,
                    bbox=dict(
                        boxstyle="round,pad=0.12",
                        facecolor="white",
                        alpha=0.65,
                        edgecolor="none"
                    ),
                    zorder=6
                )



        # Mục tiêu Ngày
        if "Mục tiêu" in df_daily.columns:
            h_mt_day, = ax_right.plot(
                df_daily["X"],
                df_daily["Mục tiêu"].replace(0, np.nan),
                color=CHART_THEME["line"]["target_day"],
                linestyle=CHART_THEME["linestyle"]["target"],
                marker=CHART_THEME["marker"]["target"],
                linewidth=1.8
            )

        # Còn lại Ca 1
        if "Lũy kế Ca1" in df_daily.columns:
            h_lk_ca1, = ax_right.plot(
                df_daily["X"],
                df_daily["Lũy kế Ca1"].replace(0, np.nan),
                color=CHART_THEME["line"]["remain_ca1"],
                linestyle=CHART_THEME["linestyle"]["remain"],
                marker=CHART_THEME["marker"]["remain"],
                markersize=5,
                linewidth=1.4
            )

        # Còn lại Ca 2
        if "Lũy kế Ca2" in df_daily.columns:
            h_lk_ca2, = ax_right.plot(
                df_daily["X"],
                df_daily["Lũy kế Ca2"].replace(0, np.nan),
                color=CHART_THEME["line"]["remain_ca2"],
                linestyle=CHART_THEME["linestyle"]["remain"],
                marker=CHART_THEME["marker"]["remain"],
                markersize=5,
                linewidth=1.4
            )

        vals = pd.concat([
            df_daily["Ca1"],
            df_daily["Ca2"],
            df_daily["Mục tiêu"],
            df_daily["Lũy kế Ca1"],
            df_daily["Lũy kế Ca2"]
        ], axis=0).replace([np.inf, -np.inf], np.nan)

        right_max = vals.dropna().max()

        if pd.isna(right_max) or right_max <= 0:
            right_max = 1

        ax_right.set_ylim(0, right_max * RIGHT_Y_MULT)

    # ======================
    # CATEGORY X
    # ======================
    labels = df_month["Label"].tolist() + df_daily["Label"].tolist()
    positions = np.arange(len(labels))

    STEP = 2
    ax_left.xaxis.set_major_locator(FixedLocator(positions[::STEP]))
    ax_left.xaxis.set_major_formatter(FixedFormatter(labels[::STEP]))
    ax_left.set_xlim(-0.5, len(labels) - 0.5)

    # ======================
    # LEGEND – GIỮ NGUYÊN NGHIỆP VỤ
    # ======================
    legend_handles = [
        h_mt_month,
        h_ca1_month,
        h_ca2_month,
        h_mt_day,
        h_lk_ca1,
        h_lk_ca2,
        h_ca1_day,
        h_ca2_day,
    ]

    legend_labels = [
        f"Mục tiêu {unit} h/thành (Tháng)",
        f"{unit} h/thành Ca 1 (Tháng)",
        f"{unit} h/thành Ca 2 (Tháng)",
        f"Mục tiêu {unit} h/thành (Ngày)",
        f"Mục tiêu {unit} Ca 1 còn lại (Ngày)",
        f"Mục tiêu {unit} Ca 2 còn lại (Ngày)",
        f"{unit} h/thành Ca 1 (Ngày)",
        f"{unit} h/thành Ca 2 (Ngày)",
    ]

    ax_left.legend(
        legend_handles,
        legend_labels,
        loc="upper center",
        bbox_to_anchor=(0.5, -0.125),
        ncol=4,
        frameon=False,
        fontsize=9
        # columnspacing=CHART_THEME["legend"]["columnspacing"],
        # handletextpad=CHART_THEME["legend"]["handletextpad"]
    )

    # ======================
    # AXIS & GRID
    # ======================
    ax_left.grid(
        axis="y",
        linestyle=CHART_THEME["grid"]["linestyle"],
        alpha=CHART_THEME["grid"]["alpha"]
    )

    ax_left.set_ylabel("SL (Tháng)", **CHART_THEME["ylabel"])
    ax_left.yaxis.set_label_coords(-0.005, 1.02)

    ax_right.set_ylabel("SL (Ngày)", **CHART_THEME["ylabel"])
    ax_right.yaxis.set_label_coords(1.005, 1.10)
    from matplotlib.ticker import MaxNLocator

    # Trục Y trái – SL Tháng
    ax_left.tick_params(axis="y", labelsize=9)
    ax_left.yaxis.set_major_locator(MaxNLocator(4))

    # Trục Y phải – SL Ngày
    ax_right.tick_params(axis="y", labelsize=9)
    ax_right.yaxis.set_major_locator(MaxNLocator(4))

    # Trục X – Tháng / Ngày
    ax_left.tick_params(axis="x", labelsize=9)
    ax_left.set_xlim(-0.65, len(labels) - 0.5)


    # ======================
    # TITLE
    # ======================
    fig.suptitle(
        title,
        fontsize=14,
        y=0.98
    )

    # ======================
    # LAYOUT & SAVE
    # ======================
    plt.subplots_adjust(
        left=0.035,
        right=0.965,
        top=0.88,
        bottom=0.24
    )
    save_figure(fig, output_file)
    
def draw_costdown_noitac_chart_flat_axis(df_mt, df_tt, title, output_file):
    """
    Costdown nội tác – chuẩn dashboard (CÁCH A):
    • MT / TT trên trục X chính
    • Tháng (T1, T2, ...) trên trục X phụ
    • MT–TT canh đối xứng quanh tâm tháng
    """

    import numpy as np
    import matplotlib.pyplot as plt
    from matplotlib.patches import Patch
    from matplotlib.lines import Line2D

    months = df_mt["Tháng"].tolist()
    num_months = len(months)

    # =========================
    # TRỤC THÁNG (TÂM)
    # =========================
    x_center = np.arange(num_months)

    pair_offset = 0.18
    bar_width  = 0.32

    x_mt = x_center - pair_offset
    x_tt = x_center + pair_offset

    # =========================
    # FIGURE – TV 65”
    # =========================
    fig, ax = plt.subplots(figsize=(13,3.5))
    ax2 = ax.twinx()

    # =========================
    # MÀU THEO LOẠI NỘI TÁC
    # =========================
    noi_tac_cfg = [
        ("NoiTac_1", "Nội tác LK tồn kho", "#9ad9a1"),
        ("NoiTac_2", "Nội tác Khuôn thành hình", "#f4b6b6"),
        ("NoiTac_3", "Nội tác Khuôn dập", "#a8c8f0"),
    ]

    # =========================
    # STACKED BAR – MT & TT
    # =========================
    for i in range(num_months):

        # ----- MT -----
        bottom = 0
        for col, _, color in noi_tac_cfg:
            val = df_mt[col].iloc[i]
            if val > 0:
                ax.bar(x_mt[i], val, bottom=bottom, width=bar_width, color=color)
                bottom += val

        if bottom > 0:
            ax.text(x_mt[i], bottom * 1.01, int(bottom),
                    ha="center", va="bottom", fontsize=10)

        # ----- TT -----
        bottom = 0
        for col, _, color in noi_tac_cfg:
            val = df_tt[col].iloc[i]
            if val > 0:
                ax.bar(x_tt[i], val, bottom=bottom, width=bar_width, color=color)
                bottom += val

        if bottom > 0:
            ax.text(x_tt[i], bottom * 1.01, int(bottom),
                    ha="center", va="bottom", fontsize=10)

    
    

    y2_max = max(df_mt["MT_LuyKe"].max(), df_tt["TT_LuyKe"].max())
    # Offset cho TT lũy kế (Triệu VND)
    # tt_offset = y2_max * 0.15   # ~15% trục Y phải (đẹp cho TV)
    label_offset = y2_max * 0.03   # ~ 3%
    # =========================
    # LINE – LŨY KẾ (Ở TÂM THÁNG)
    # =========================
    ax2.plot(
        x_center,
        df_mt["MT_LuyKe"],
        color="red",
        marker="o",
        linewidth=2.1,
        label="MT lũy kế"
    )

    ax2.plot(
        x_center,
        df_tt["TT_LuyKe"],
        color="#1f78b4",
        marker="o",
        linewidth=2.1,
        label="TT lũy kế"
    )



    for i in range(num_months):
        if df_mt["MT_LuyKe"].iloc[i] > 0:
            ax2.text(x_center[i], df_mt["MT_LuyKe"].iloc[i] + label_offset,
                     int(df_mt["MT_LuyKe"].iloc[i]),
                     color="red", fontsize=10,
                     ha="center", va="bottom")

        if df_tt["TT_LuyKe"].iloc[i] > 0:
            ax2.text(x_center[i], df_tt["TT_LuyKe"].iloc[i] -  label_offset,
                     int(df_tt["TT_LuyKe"].iloc[i]),
                     color="#1f78b4", fontsize=10,
                     ha="center", va="top")

    # =================================================
    # TRỤC X – DÒNG 1: MT / TT
    # =================================================
    xticks_main = np.ravel(np.column_stack([x_mt, x_tt]))
    ax.set_xticks(xticks_main)
    ax.set_xticklabels(
        ["MT", "TT"] * num_months,
        fontsize=11
    )

    # =================================================
    # TRỤC X PHỤ – DÒNG 2: THÁNG (T1, T2, ...)
    # =================================================
    ax_month = ax.secondary_xaxis("bottom")
    ax_month.set_xticks(x_center)
    ax_month.set_xticklabels(
        [f"T{i+1}" for i in range(num_months)],
        fontsize=11,
        fontweight="bold"
    )

    ax_month.spines["bottom"].set_position(("outward", 18))
    ax_month.spines["bottom"].set_visible(False)
    ax_month.tick_params(length=0)

    # =========================
    # TITLE & AXIS
    # =========================
    fig.suptitle(
        title,
        fontsize=18,
        # fontweight="bold",
        y=1.02
    )

    ax.set_ylabel("Đơn tháng\n(Triệu VND)", rotation=0, fontsize=11)
    ax.yaxis.set_label_coords(-0.019, 1.01)

    ax2.set_ylabel("Lũy kế\n(Triệu VND)", rotation=0, fontsize=11)
    ax2.yaxis.set_label_coords(1.02, 1.13)


    ax.grid(axis="y", linestyle="--", alpha=0.25)
    ax.set_axisbelow(True)
    ax.spines["top"].set_visible(False)
    ax2.spines["top"].set_visible(False)

    # =========================
    # LEGEND
    # =========================
    legend_elements = [
        Patch(facecolor="#9ad9a1", label="Nội tác LK tồn kho"),
        Patch(facecolor="#f4b6b6", label="Nội tác Khuôn thành hình"),
        Patch(facecolor="#a8c8f0", label="Nội tác Khuôn dập"),
        Line2D([0], [0], color="red", lw=2.5, marker="o", label="MT lũy kế"),
        Line2D([0], [0], color="#1f78b4", lw=2.5, marker="o", label="TT lũy kế"),
    ]

    ax.legend(
        handles=legend_elements,
        loc="upper center",
        bbox_to_anchor=(0.5, -0.12),
        ncol=5,
        frameon=False,
        fontsize=11
    )

    # =========================
    # SAVE – TV 65”
    # =========================
    plt.subplots_adjust(left=0.06, right=0.96, top=0.90, bottom=0.1)
    plt.savefig(
        output_file.with_suffix(".png"),
        dpi=240,
        bbox_inches="tight",
        pad_inches=0.04
    )
    plt.close()

def draw_costdown_mt_tt_month_chart(df, title, output_file):
    x = np.arange(len(df))
    width = 0.35

    # fig, ax = plt.subplots(figsize=(12, 4))
    fig, ax = plt.subplots(figsize=(32/2.54, 9/2.54))   # ≈ (12.6, 3.54)


    bars_mt = ax.bar(
        x - width/2,
        df["Mục tiêu"],
        width,
        color="red",
        label="Mục tiêu"
    )

    bars_tt = ax.bar(
        x + width/2,
        df["Thực tích"],
        width,
        color="#2f80ed",
        label="Thực tích"
    )

    # Data label
    for bars in (bars_mt, bars_tt):
        for b in bars:
            h = b.get_height()
            if h > 0:
                ax.text(
                    b.get_x() + b.get_width()/2,
                    h,
                    f"{int(h)}",
                    ha="center",
                    va="bottom",
                    fontsize=9
                )

    ax.set_xticks(x)
    ax.set_xticklabels(df["Hạng mục"], rotation=0)

    ax.set_ylabel("Triệu\n(VND)", rotation=0, fontsize=11)
    ax.yaxis.set_label_coords(-0.01, 1.02)

    fig.suptitle(
        title,
        fontsize=18,
        # fontweight="bold",
        # color="blue",
        y=1
    )

    ax.grid(axis="y", linestyle="--", alpha=0.25)
    ax.legend(
        loc="upper center",
        bbox_to_anchor=(0.5, -0.05),
        ncol=5,
        frameon=False,
        fontsize=11
    )

    plt.subplots_adjust(
        left=0.04,
        right=0.96,
        top=0.90,
        bottom=0.1
    )
    plt.savefig(output_file.with_suffix(".png"),dpi=240,bbox_inches="tight",pad_inches=0.04)
    plt.close()



def draw_timer_utilization_chart(df, title, output_file,show_title=True):
    x = np.arange(len(df)) *1.5

    # ======================
    # FIGURE – FULL SLIDE
    # ======================
    fig, ax = plt.subplots(figsize=(30/2.54, 16/2.54))

    # ======================
    # DATA (%)
    # ======================
    co_tai = df["% Có tải"] * 100
    khong_tai = df["% Không tải"] * 100
    tat_may = df["% Tắt máy"] * 100
    target = df["% Mục tiêu"] * 100

    # ======================
    # STACKED BAR – NHÂN VẬT CHÍNH
    # ======================
    BAR_WIDTH = 0.2
    bars_co = ax.bar(
        x,
        co_tai,
        color="#9ACD32",
        label="% Có tải"
    )

    ax.bar(
        x,
        khong_tai,
        bottom=co_tai,
        color="#FFD966",
        label="% Không tải"
    )

    ax.bar(
        x,
        tat_may,
        bottom=co_tai + khong_tai,
        color="#A6A6A6",
        label="% Tắt máy"
    )

    # ======================
    # DATA LABEL – CHỈ % CÓ TẢI
    # ======================
    for bar, v in zip(bars_co, co_tai):
        if v >= 0:   # ✅ chỉ hiện khi đủ đáng chú ý
            ax.text(
                bar.get_x() + bar.get_width()/2,
                v / 2,
                f"{int(round(v))}%",
                ha="center",
                va="center",
                fontsize=11,
                fontweight="bold",
                color="black",
                bbox=dict(
                    boxstyle="round,pad=0.2",
                    facecolor="white",
                    alpha=0.75,
                    edgecolor="none"
                )
            )

    # ======================
    # LINE – % MỤC TIÊU (PHỤ)
    # ======================
    ax.plot(
        x,
        target,
        color="red",
        linestyle="--",
        marker="o",
        markersize=6,
        linewidth=2.5,   # ✅ dày hơn chút
        label="% Mục tiêu",
        zorder=5
    )

    # DATA LABEL – MỤC TIÊU (TRÊN LINE)
    for i, v in enumerate(target):
        if v > 0:
            ax.text(
                x[i],
                v + 2.5,
                f"{int(round(v))}%",
                ha="center",
                va="bottom",
                fontsize=10,
                color="red",
                bbox=dict(
                    boxstyle="round,pad=0.15",
                    facecolor="white",
                    alpha=0.9,
                    edgecolor="none"
                )
            )


    # ======================
    # AXIS
    # ======================
    ax.set_xticks(x)
    ax.set_xticklabels(df["Hạng mục"], fontsize=11)
    ax.set_ylim(0, 100)

    ax.set_ylabel("Tỷ lệ (%)", **CHART_THEME["ylabel"])
    ax.yaxis.set_label_coords(-0.01, 1.02)

    # ======================
    # GRID & STYLE
    # ======================
    ax.grid(
        axis="y",
        linestyle=CHART_THEME["grid"]["linestyle"],
        alpha=0.22
    )
    ax.set_axisbelow(True)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#AAAAAA")
    ax.spines["bottom"].set_color("#AAAAAA")

    # ======================
    # TITLE
    # ======================
    if show_title and title:
        fig.suptitle(
            title,
            fontsize=18,
            fontweight="bold",
            y=0.98
        )

    # ======================
    # LEGEND
    # ======================
    ax.legend(
        loc="lower center",
        bbox_to_anchor=(0.5, -0.12),
        ncol=4,
        frameon=False,
        fontsize=CHART_THEME["legend"]["fontsize"]
    )

    # ======================
    # LAYOUT & SAVE
    # ======================
    plt.subplots_adjust(
        left=0.055,
        right=1,
        top=0.95,
        bottom=0.11
    )

    save_figure(fig, output_file)

def draw_timer_block_chart_manual_style(df, title, output_file):
    x = np.arange(len(df)) *1.15
    STEP = max(1, len(df) // 15)   # ~10 label là đẹp
    # fig, ax = plt.subplots(figsize=(15, 4.8))
    fig, ax = plt.subplots(figsize=(13,3))   # ≈ (12.6, 3.54)

    bar_width = 0.35

    # stacked bar (0–1 → %)
    co_tai = df["% Có tải"] * 100
    khong_tai = df["% Không tải"] * 100
    tat_may = df["% Tắt máy"] * 100

    ax.bar(x, co_tai, width=bar_width, color="#9ACD32", label="% Có tải")
    ax.bar(x, khong_tai, width=bar_width,  bottom=co_tai, color="#FFD966", label="% Không tải")
    ax.bar(x, tat_may, width=bar_width,  bottom=co_tai + khong_tai, color="#A6A6A6", label="% Tắt máy")


    SHOW_TEXT = len(df) <= 30

    if SHOW_TEXT:
        for i, v in enumerate(co_tai):
            if v > 5:
                ax.text(
                    x[i],
                    v / 2,
                    f"{int(v)}",
                    ha="center",
                    va="center",
                    fontsize=7,
                    fontweight="bold",
                    color="black",
                    bbox=dict(
                        boxstyle="round,pad=0.18",
                        facecolor="white",
                        alpha=0.75,
                        edgecolor="none"
                    )
                )

    target_c1, target_c2 = get_targets_from_df(df)

    if target_c1 is not None:
        ax.plot(
            x, [target_c1]*len(x),
            color="#9B5DE5",
            linestyle="--",
            linewidth=2.2,
            marker="o",
            markersize=4,
            label="Mục tiêu Ca 1",
            zorder=6
        )

    if target_c2 is not None:
        ax.plot(
            x, [target_c2]*len(x),
            color="#D62828",
            linestyle="--",
            linewidth=2.2,
            marker="s",
            markersize=4,
            label="Mục tiêu Ca 2",
            zorder=6
        )


    ax.set_ylim(0, 100)
    ax.set_xticks(x[::STEP])
    ax.set_xticklabels(df["Thời gian"].iloc[::STEP], fontsize=9)
    ax.set_ylabel("Tỷ lệ (%)", rotation=0, fontsize=9)
    ax.yaxis.set_label_coords(-0.005, 1.042)
    fig.suptitle(title, fontsize=14, y=1.09)
    # ax.set_ylabel("Tỷ lệ (%)")
    # ax.set_title(title, fontsize=14, fontweight="bold")

    ax.grid(axis="y", linestyle="--", alpha=0.18)
    ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.1), ncol=5, frameon=False, fontsize=9,columnspacing=1.2)

    plt.subplots_adjust(left=0.06, right=0.96, top=0.95, bottom=0.1)
    plt.savefig(output_file.with_suffix(".png"),dpi=240,bbox_inches="tight",pad_inches=0.04)
    plt.close()

# ========== PPT ==========
def replace_image_in_ppt(prs, shape_name, image_path, slide_index):
    slide = prs.slides[slide_index]

    print(f"🔎 Tìm anchor '{shape_name}' trong slide {slide_index}")

    found = False

    for shape in slide.shapes:
        print("   └─ shape.name =", repr(shape.name))

        if shape.name == shape_name:
            found = True

            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height

            slide.shapes._spTree.remove(shape._element)

            pic = slide.shapes.add_picture(
                str(image_path),
                left,
                top,
                width=width,
                height=height
            )
            pic.name = shape_name

            # prs.save(PPT_FILE)
            print(f"✅ Đã update ảnh tại anchor {shape_name}")
            return

    if not found:
        print(f"❌ KHÔNG tìm thấy anchor '{shape_name}' trong slide {slide_index}")  
         


def update_title_in_ppt(
    prs,
    slide_index,
    title_text,
    title_shape_name="title_block",
    font_name="Calibri",
    font_size_pt=20,
    bold=True,
    color_rgb=(0, 0, 0)
):
    slide = prs.slides[slide_index]

    for shape in slide.shapes:
        if shape.name == title_shape_name and shape.has_text_frame:
            tf = shape.text_frame
            tf.clear()

            p = tf.paragraphs[0]
            p.text = title_text
            p.font.name = font_name
            p.font.size = Pt(font_size_pt)
            p.font.bold = bold
            p.font.italic = False
            p.font.color.rgb = RGBColor(*color_rgb)
            return

    print(f"⚠️ Không tìm thấy title '{title_shape_name}' trong slide {slide_index}")

# =================================================
# YEAR SUMMARY – BV + PCS COMBO (A1:I13)
# =================================================
def load_year_summary():
    df = pd.read_excel(
        EXCEL_YEAR_FILE,
        sheet_name=EXCEL_YEAR_SHEET,
        usecols="A:I",
        nrows=13
    )

    df.columns = [
        "Tháng",
        "Ngày làm việc",
        "MT Giờ",
        "MT BV",
        "TT BV",
        "P/MH BV",
        "MT PCS",
        "TT PCS",
        "P/MH PCS"
    ]

    df = df[df["Tháng"].notna()].reset_index(drop=True)
    # ✅ CHUẨN HOÁ: 1月 → Tháng 1
    df["Tháng"] = df["Tháng"].astype(str).str.replace("月", "", regex=False)
    df["Tháng"] = "Tháng " + df["Tháng"]

    return df


def draw_year_bv_pcs_combo_chart(df, title, output_file,show_title=True):
    x = np.arange(len(df))
    width = 0.32

    # ======================
    # FIGURE – FULL SLIDE PPT
    # ======================
    fig, ax = plt.subplots(figsize=(30/2.54, 16/2.54))
    ax2 = ax.twinx()

    # ======================
    # COLOR SYSTEM
    # ======================
    COL_BV_TARGET = "#B7D9F7"   # xanh nhạt – mục tiêu
    COL_BV_ACTUAL = "#1F6FB2"   # xanh đậm – thực tích
    COL_PCS_TARGET = "#F39C12"  # cam – mục tiêu
    COL_PCS_ACTUAL = "#7D3C98"  # tím trầm – thực tích

    # ======================
    # BAR – BẢN VẼ (NHÂN VẬT CHÍNH)
    # ======================
    bars_mt_bv = ax.bar(
        x - width/2,
        df["MT BV"],
        width=width,
        color=COL_BV_TARGET,
        label="Mục tiêu Bản vẽ"
    )

    bars_tt_bv = ax.bar(
        x + width/2,
        df["TT BV"],
        width=width,
        color=COL_BV_ACTUAL,
        label="Thực tích Bản vẽ"
    )

    # ======================
    # DATA LABEL – BAR BẢN VẼ (CENTER + WHITE HALO)
    # ======================
    BAR_TEXT_BOX = dict(
        boxstyle="round,pad=0.18",
        facecolor="white",
        alpha=0.75,
        edgecolor="none"
    )

    for bars in (bars_mt_bv, bars_tt_bv):
        for b in bars:
            h = b.get_height()
            if h > 0:
                ax.text(
                    b.get_x() + b.get_width() / 2,
                    h / 2,                     # ✅ nằm chính giữa cột
                    f"{int(h)}",
                    ha="center",
                    va="center",
                    fontsize=8,
                    color="black",
                    bbox=BAR_TEXT_BOX
                )

    # ======================
    # LINE – PCS (NHÂN VẬT PHỤ)
    # ======================
    ax2.plot(
        x,
        df["MT PCS"],
        color=COL_PCS_TARGET,
        linestyle="-",
        marker="o",
        linewidth=2.0,
        label="Mục tiêu PCS"
    )

    ax2.plot(
        x,
        df["TT PCS"],
        color=COL_PCS_ACTUAL,
        linestyle="--",     # ✅ NÉT ĐỨT = PHỤ
        marker=None,        # ✅ KHÔNG MARKER
        linewidth=1.6,
        alpha=0.85,
        label="Thực tích PCS"
    )

    # ======================
    # DATA LABEL – PCS (HALO + OFFSET ĐỀU)
    # ======================
    PCS_TEXT_BOX = dict(
        boxstyle="round,pad=0.18",
        facecolor="white",
        alpha=0.85,
        edgecolor="none"
    )

    y_pcs_max = ax2.get_ylim()[1]
    offset = y_pcs_max * 0.02   # tăng nhẹ so với 0.015

    # ----- MT PCS -----
    for i, y in enumerate(df["MT PCS"]):
        if y > 0:
            ax2.text(
                x[i],
                y + offset,
                f"{int(y)}",
                ha="center",
                va="bottom",
                fontsize=8,
                color=COL_PCS_TARGET,
                bbox=PCS_TEXT_BOX
            )

    # ----- TT PCS -----
    for i, y in enumerate(df["TT PCS"]):
        if y > 0:
            ax2.text(
                x[i],
                y - offset,
                f"{int(y)}",
                ha="center",
                va="top",
                fontsize=8,
                color=COL_PCS_ACTUAL,
                bbox=PCS_TEXT_BOX
            )


    # ======================
    # AXIS
    # ======================
    ax.set_xticks(x)
    ax.set_xticklabels(df["Tháng"], fontsize=9)

    ax.set_ylabel("Số Bản vẽ", **CHART_THEME["ylabel"])
    ax.yaxis.set_label_coords(-0.01, 1.02)

    ax2.set_ylabel("Số PCS", **CHART_THEME["ylabel"])
    ax2.yaxis.set_label_coords(1.03, 1.025)

    # SCALE Y – RẤT QUAN TRỌNG CHO CÂN MẮT
    y_bv = pd.concat([df["MT BV"], df["TT BV"]]).dropna()
    y_max_bv = y_bv.max() if not y_bv.empty else 1
    ax.set_ylim(0, y_max_bv * 1.1)

    y_pcs = pd.concat([df["MT PCS"], df["TT PCS"]]).dropna()
    y_max_pcs = y_pcs.max() if not y_pcs.empty else 1
    ax2.set_ylim(0, y_max_pcs * 1.15)

    # ======================
    # GRID & STYLE
    # ======================
    ax.grid(
        axis="y",
        linestyle=CHART_THEME["grid"]["linestyle"],
        alpha=0.2
    )
    ax.set_axisbelow(True)
    ax.spines["top"].set_visible(False)
    ax2.spines["top"].set_visible(False)

    # ======================
    # TITLE – DỄ CHỊU CHO 1 SLIDE
    # ======================
    if show_title and title:
        fig.suptitle(
            title,
            fontsize=18,
            fontweight="bold",
            y=0.96
        )

    # ======================
    # LEGEND – GỌN, KHÔNG ĐÈ
    # ======================
    h1, l1 = ax.get_legend_handles_labels()
    h2, l2 = ax2.get_legend_handles_labels()

    ax.legend(
        h1 + h2,
        l1 + l2,
        loc="lower center",
        bbox_to_anchor=(0.5, -0.12),
        ncol=4,
        frameon=False,
        fontsize=CHART_THEME["legend"]["fontsize"]
    )

    # ======================
    # LAYOUT – CÂN ĐỐI SLIDE
    # ======================
    plt.subplots_adjust(
        left=0.045,
        right=0.94,
        top=0.95,
        bottom=0.11
    )

    # ======================
    # SAVE
    # ======================
    save_figure(fig, output_file)


def process_year_summary_combo(prs):
    df = load_year_summary()

    img_path_local = LOCAL_OUT_DIR / "year_summary_bv_pcs_combo.png"

    draw_year_bv_pcs_combo_chart(df,"TỔNG BẢN VẼ & PCS HOÀN THÀNH THEO THÁNG – NĂM 2026ff",img_path_local,show_title=False)# ⬅️ TẮT TITLE

    slide_idx = PPT_ANCHORS["year"]["slide"]

    # ✅ THÊM DÒNG NÀY
    update_title_in_ppt(
        prs,
        slide_idx,
        "TỔNG BẢN VẼ & PCS HOÀN THÀNH THEO THÁNG – NĂM 2026",
        title_shape_name="title_year",
        font_name="Calibri",
        font_size_pt=20,
        bold=True,
        color_rgb=(0, 0, 128)
    )

    
    # ✅ COPY LÊN SERVER
    copy_to_server_safe(img_path_local, SERVER_OUT_DIR)


    replace_image_in_ppt(
        prs,
        PPT_ANCHORS["year"]["year_tsx"],
        img_path_local,
        slide_index=slide_idx
    )

def load_file3_G1_L16():
    raw = pd.read_excel(
        EXCEL_THUCTICH_FILE,
        sheet_name=0,
        usecols="G:L",
        header=None
    )

    # ✅ CỘT ĐẦU TIÊN = iloc[:, 0]
    first_col = (
        raw.iloc[:, 0]
        .astype(str)
        .str.replace("\xa0", " ", regex=False)
        .str.strip()
        .str.lower()
    )

    # ✅ TÌM DÒNG HEADER "Công đoạn"
    header_idx = first_col[first_col == "công đoạn"].index

    if header_idx.empty:
        raise ValueError("❌ Không tìm thấy dòng 'Công đoạn' trong file Excel")

    header_row = header_idx[0]

    # ✅ LẤY DATA PHÍA DƯỚI HEADER
    df = raw.iloc[header_row + 1 :].copy()

    df.columns = [
        "Công đoạn",
        "Phụ tải CĐ hiện tại",
        "Phụ tải CĐ tổng",
        "Số bản vẽ trễ hẹn",
        "Số bản vẽ hoàn thành",
        "Hiệu suất (Ngày)"
    ]

    # ✅ CLEAN
    df = df[df["Công đoạn"].notna()].reset_index(drop=True)
    df["Công đoạn"] = df["Công đoạn"].astype(str).str.strip()

    return df

def load_costdown_mt_tt_month_from_matrix():
    """
    Đọc bảng AA4:AD7 (dạng ma trận) và chuẩn hóa để vẽ biểu đồ
    """
    df_raw = pd.read_excel(
        EXCEL_COSTDOWN_FILE,
        sheet_name=EXCEL_COSTDOWN_SHEET,
        usecols="AA:AD",
        skiprows=4,   # bắt đầu từ row 5
        nrows=3       # row 5,6,7
    )

    # Đặt tên cột
    df_raw.columns = [
        "Loại",           # AA
        "Nội tác tồn kho",
        "Nội tác khuôn thành hình",
        "Nội tác khuôn dập"
    ]

    # Lấy dòng Mục tiêu & Thực tích
    df_mt = df_raw[df_raw["Loại"].str.contains("Mục tiêu", na=False)].iloc[0]
    df_tt = df_raw[df_raw["Loại"].str.contains("Thực tích", na=False)].iloc[0]

    # Tạo dataframe chuẩn để vẽ
    df = pd.DataFrame({
        "Hạng mục": [
            "1.Nội tác linh kiện tồn kho",
            "2.Nội tác linh kiện khuôn thành hình",
            "3.Nội tác linh kiện khuôn dập"
        ],
        "Mục tiêu": [
            df_mt["Nội tác tồn kho"],
            df_mt["Nội tác khuôn thành hình"],
            df_mt["Nội tác khuôn dập"]
        ],
        "Thực tích": [
            df_tt["Nội tác tồn kho"],
            df_tt["Nội tác khuôn thành hình"],
            df_tt["Nội tác khuôn dập"]
        ]
    })

    return df

def load_tong_hop_timer_data():
    raw = pd.read_excel(
        EXCEL_TIMER_FILE,
        sheet_name=EXCEL_TIMER_SHEET,
        header=None,
        dtype=str
    )

    # ===== CHUẨN HÓA CỘT A =====
    raw[0] = (
        raw[0]
        .ffill()
        .astype(str)
        .str.replace("\xa0", " ", regex=False)
        .str.strip()
        .str.lower()
    )

    def find_row_regex(pattern):
        idx = raw[0].str.match(pattern, na=False)
        if not idx.any():
            raise ValueError(f"❌ Không tìm thấy dòng regex: {pattern}")
        return idx.idxmax()

    # ✅ MATCH CHÍNH XÁC – KHÔNG DÍNH TIÊU ĐỀ
    row_hang_muc  = find_row_regex(r"^hạng mục$")
    row_muc_tieu  = find_row_regex(r"^% *mục tiêu$")
    row_co_tai    = find_row_regex(r"^% *có tải$")
    row_khong_tai = find_row_regex(r"^% *không tải$")
    row_tat_may   = find_row_regex(r"^% *tắt máy$")

    # ✅ CÔNG ĐOẠN
    cong_doan = raw.loc[row_hang_muc, 1:].dropna().str.strip()

    def parse_percent(row_idx):
        return (
            raw.loc[row_idx, 1:1+len(cong_doan)]
            .astype(str)
            .str.replace("%", "", regex=False)
            .str.replace(",", ".", regex=False)
            .astype(float)
        )

    df = pd.DataFrame({
        "Hạng mục": cong_doan.values,
        "% Có tải": parse_percent(row_co_tai),
        "% Không tải": parse_percent(row_khong_tai),
        "% Tắt máy": parse_percent(row_tat_may),
        "% Mục tiêu": parse_percent(row_muc_tieu),
    })

    return df

def load_timer_block_sheet(sheet_name):
    """
    Đọc vùng H1:M75 của 1 sheet timer
    """
    df = pd.read_excel(
        EXCEL_TIMER_FILE,
        sheet_name=sheet_name,
        usecols="H:M",
        nrows=75
    )

    df.columns = ["Time", "CoTai", "KhongTai", "TatMay", "Tong", "GhiChu"]

    # Ép numeric an toàn
    for c in ["CoTai", "KhongTai", "TatMay"]:
        df[c] = pd.to_numeric(df[c], errors="coerce").fillna(0)

    df = df.dropna(how="all")

    df["X"] = range(len(df))
    return df

def replace_table_by_anchor_centered(
    prs,
    slide_index,
    anchor_name,
    df,
    report_date
):
    slide = prs.slides[slide_index]

    for shape in slide.shapes:
        if shape.name == anchor_name:

            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            slide.shapes._spTree.remove(shape._element)

            cols = df.shape[1]
            data_rows = len(df)

            table_shape = slide.shapes.add_table(
                # data_rows + 2,   # title + header + data
                data_rows + 1,   # header + data
                cols,
                left, top, width, height
            )
            table_shape.name = anchor_name
            table = table_shape.table

            # ===== COLUMN WIDTH =====
            for c in range(cols):
                table.columns[c].width = int(width / cols)

            # ==================================================
            # TITLE
            # ==================================================
            # title_cell = table.cell(0, 0)
            # title_cell.merge(table.cell(0, cols - 1))
            # title_cell.text = (
            #     f"Bản vẽ Hoàn thành, Trễ hẹn và Phụ tải Công đoạn – {report_date}"
            # )
            # title_cell.fill.solid()
            # title_cell.fill.fore_color.rgb = RGBColor(21, 96, 130)

            # for p in title_cell.text_frame.paragraphs:
            #     p.alignment = PP_ALIGN.CENTER
            #     p.font.size = Pt(16)
            #     p.font.bold = True
            #     p.font.color.rgb = RGBColor(255, 255, 255)

            # ==================================================
            # HEADER
            # ==================================================
            for c, col_name in enumerate(df.columns):
                cell = table.cell(0, c)
                cell.text = col_name
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(21, 96, 130)

                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.CENTER
                    p.font.size = Pt(15)
                    p.font.bold = True

            # ==================================================
            # DATA
            # ==================================================
            for r in range(data_rows):
                for c in range(cols):
                    val = df.iloc[r, c]

                    if isinstance(val, float):
                        txt = f"{val:.2f}"
                    else:
                        txt = str(val)

                    cell = table.cell(r + 1, c)
                    cell.text = txt
                    cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

                    for p in cell.text_frame.paragraphs:
                        p.alignment = PP_ALIGN.CENTER
                        p.font.size = Pt(15)

                        if df.columns[c] == "Số bản vẽ trễ hẹn":
                            p.font.color.rgb = RGBColor(255, 0, 0)

                    cell.fill.solid()
                    cell.fill.fore_color.rgb = (
                        RGBColor(226, 234, 248)
                        if r % 2 == 0
                        else RGBColor(255, 255, 255)
                    )

            print("✅ Table tạo GIỐNG HỆT HÌNH MẪU")
            return

    print(f"❌ Không tìm thấy anchor '{anchor_name}'")

def update_table_title_by_anchor(
    prs,
    slide_index,
    anchor_name,
    title_text,
    font_name="Calibri",
    font_size_pt=14,
    bold=True,
    color_rgb=(0, 0, 128)
):
    slide = prs.slides[slide_index]

    for shape in slide.shapes:
        if shape.name == anchor_name and shape.has_text_frame:
            tf = shape.text_frame
            tf.clear()

            p = tf.paragraphs[0]
            p.text = title_text
            p.font.name = font_name
            p.font.size = Pt(font_size_pt)
            p.font.bold = bold
            p.font.color.rgb = RGBColor(*color_rgb)

            print(f"✅ Đã update tiêu đề table tại '{anchor_name}'")
            return

    print(f"❌ Không tìm thấy anchor tiêu đề '{anchor_name}'") 

def build_costdown_data_noitac():
    """
    Đọc bảng Costdown nội tác đúng theo cấu trúc Excel (B:J)
    """
    df = pd.read_excel(
        EXCEL_COSTDOWN_FILE,
        sheet_name=EXCEL_COSTDOWN_SHEET,
        usecols="B:J",     # ✅ CHUẨN THEO FILE
        skiprows=2,        # bắt đầu từ dòng tiêu đề
        nrows=24
    )

    df.columns = [
        "Tháng",          # B
        "MT_TT",          # C
        "NoiTac_1",       # D
        "NoiTac_2",       # E
        "NoiTac_3",       # F
        "MT_Thang",       # G
        "TT_Thang",       # H
        "MT_LuyKe",       # I
        "TT_LuyKe"        # J
    ]

    df["Tháng"] = df["Tháng"].astype(str).str.strip()
    df["MT_TT"] = df["MT_TT"].astype(str).str.strip()

    df_mt = df[df["MT_TT"] == "Mục tiêu"].reset_index(drop=True)
    df_tt = df[df["MT_TT"] == "Thực tích"].reset_index(drop=True)

    return df_mt, df_tt

def process_costdown_noitac(prs):
    df_mt, df_tt = build_costdown_data_noitac()

    img_path_local = LOCAL_OUT_DIR / "costdown_noi_tac.png"

    draw_costdown_noitac_chart_flat_axis(
    df_mt,
    df_tt,
    "COSTDOWN NỘI TÁC – MỤC TIÊU & THỰC TÍCH NĂM 2026",
    img_path_local
)
    
    # ✅ COPY LÊN SERVER
    copy_to_server_safe(img_path_local, SERVER_OUT_DIR)



    replace_image_in_ppt(
        prs,
        PPT_ANCHORS["costdown"]["TT_noi_tac"],
        img_path_local,
        slide_index=PPT_ANCHORS["costdown"]["slide"]
    )


def process_costdown_mt_tt_month(prs):
    df = load_costdown_mt_tt_month_from_matrix()

    img_path_local = LOCAL_OUT_DIR / "costdown_mt_tt_thang_4.png"

    draw_costdown_mt_tt_month_chart(
        df,
        "MỤC TIÊU – THỰC TÍCH THÁNG 4/2026",
        img_path_local
    )


    # ✅ COPY LÊN SERVER
    copy_to_server_safe(img_path_local, SERVER_OUT_DIR)


    replace_image_in_ppt(
        prs,
        PPT_ANCHORS["costdown"]["TT_tong"],
        img_path_local,
        slide_index=PPT_ANCHORS["costdown"]["slide"]
    )

def process_timer_utilization_chart(prs):
    df = load_tong_hop_timer_data()
    print(df)
    # print(df.dtypes)
    print(df["% Có tải"] + df["% Không tải"] + df["% Tắt máy"])


    img_path_local = LOCAL_OUT_DIR / "timer_utilization.png"

    draw_timer_utilization_chart(
        df,
        "BIỂU ĐỒ TỈ LỆ HOẠT ĐỘNG MÁY TẠI CÁC CÔNG ĐOẠN TRONG THÁNG NĂM 2026",
        img_path_local,
        show_title=False
    )

    
    # ✅ COPY LÊN SERVER
    copy_to_server_safe(img_path_local, SERVER_OUT_DIR)


    replace_image_in_ppt(
        prs,
        "anchor_timer_utilization",
        img_path_local,
        slide_index=3
    )

def load_timer_sheet_block(sheet_name):
    df = pd.read_excel(
        EXCEL_TIMER_FILE,
        sheet_name=sheet_name,
        usecols="H:M",
        header=0
    )

    df.columns = [
        "Thời gian",
        "Mục tiêu C1",
        "Mục tiêu C2",
        "% Có tải",
        "% Không tải",
        "% Tắt máy",
    ]

    for c in ["Mục tiêu C1", "Mục tiêu C2", "% Có tải", "% Không tải", "% Tắt máy"]:
        df[c] = pd.to_numeric(df[c], errors="coerce")

    # ✅ chỉ giữ dòng có dữ liệu thực
    df = df[df["% Có tải"].notna()].reset_index(drop=True)

    return df

def get_targets_from_df(df):
    """
    Trả về mục tiêu C1/C2 dạng % (65, 80)
    """
    def normalize(v):
        if pd.isna(v) or v <= 0:
            return None
        return int(round(v * 100)) if v <= 1 else int(round(v))

    c1 = df["Mục tiêu C1"].dropna()
    c2 = df["Mục tiêu C2"].dropna()

    t1 = normalize(c1.iloc[0]) if not c1.empty else None
    t2 = normalize(c2.iloc[0]) if not c2.empty else None

    return t1, t2
def process_timer_block_charts(prs):
    """
    Vẽ biểu đồ TIMER cho từng công đoạn
    và gắn vào anchor_timer của slide block tương ứng
    """

    for sheet_name in TIMER_SHEETS:

        slide_idx = BLOCK_SLIDE_MAP.get(sheet_name)
        if slide_idx is None:
            print(f"⏭️ Bỏ qua {sheet_name} (không có slide)")
            continue

        print(f"⏳ Xử lý Timer sheet: {sheet_name}")

        # ===== LOAD DATA =====
        try:
            df = load_timer_sheet_block(sheet_name)
        except Exception as e:
            print(f"⚠️ Lỗi đọc sheet {sheet_name}: {e}")
            continue

        if df.empty:
            print(f"⚠️ Sheet {sheet_name} không có dữ liệu timer")
            continue

        # ===== LẤY MỤC TIÊU C1 / C2 =====
        target_c1, target_c2 = get_targets_from_df(df)

        mt_text = []
        if target_c1 is not None:
            # mt_text.append(f"Ca 1={int(round(target_c1*100 if target_c1 <= 1 else target_c1))}%")
            mt_text.append(f"Ca 1={target_c1}%")
        if target_c2 is not None:
            # mt_text.append(f"Ca 2={int(round(target_c2*100 if target_c2 <= 1 else target_c2))}%")
            mt_text.append(f"Ca 2={target_c2}%")

        mt_str = ", ".join(mt_text) if mt_text else "Chưa thiết lập MT"

        title = (
            f"BIỂU ĐỒ TỈ LỆ HOẠT ĐỘNG MÁY – CĐ {sheet_name} – 8.5h/ca -"
            f"[MT có tải: {mt_str}]"
        )

        img_path_local = LOCAL_OUT_DIR / f"timer_block_{sheet_name}.png"

        # ===== VẼ BIỂU ĐỒ TIMER =====
        draw_timer_block_chart_manual_style(
            df=df,
            title=title,
            output_file=img_path_local
        )

        
        # ✅ COPY LÊN SERVER
        copy_to_server_safe(img_path_local, SERVER_OUT_DIR)

        # ===== GẮN VÀO ANCHOR TIMER =====
        replace_image_in_ppt(
            prs,
            "anchor_timer",
            img_path_local,
            slide_index=slide_idx
        )

        print(f"✅ Done timer block {sheet_name}")



# ========== MAIN ==========
def main():
    # ✅ 1. LOAD EXCEL
    df_raw = load_raw_excel()

    # ✅ LOAD PPT 1 LẦN
    prs = Presentation(PPT_FILE)

    # =================================================
    # GLOBAL charts
    # =================================================
    # ✅ LẤY SLIDE GLOBAL
    slide_idx = PPT_ANCHORS["global"]["slide"]
    title_text = "MỤC TIÊU – THỰC TÍCH BẢN VẼ & SỐ PCS HOÀN THÀNH THEO NGÀY"
    update_title_in_ppt(prs, slide_idx, title_text, title_shape_name="title_global",font_name="Calibri",font_size_pt=20,bold=True,color_rgb=(0, 0, 128))   # xanh đậm
    for cfg in GLOBAL_CONFIG:
        try:
            df_chart = build_chart_df_global(df_raw, cfg["rows"])
        except Exception as e:
            print(f"⚠️ Bỏ qua GLOBAL {cfg['name']} vì lỗi: {e}")
            continue

        img_path_local = LOCAL_OUT_DIR / f"{cfg['name']}.png"
        draw_combo_chart_global(
            df=df_chart,
            title=cfg.get("title", cfg["name"]),
            ylabel="Số lượng",
            output_file=img_path_local,
            legend_prefix=cfg.get("legend_prefix", "")
        )
        


        # ✅ LẤY ANCHOR
        anchor = PPT_ANCHORS["global"][cfg["name"]]


        
        # ✅ COPY LÊN SERVER
        copy_to_server_safe(img_path_local, SERVER_OUT_DIR)
    

        # ✅ GẮN ẢNH VÀO PPT
        replace_image_in_ppt(
            prs,
            anchor,
            img_path_local,
            slide_index=slide_idx
        )

 

    # =================================================
    # BLOCK charts
    # =================================================
    blocks = split_blocks(df_raw)   # ✅ df_raw đã tồn tại

    def get_cong_doan_name(block):
        """
        Excel:
        A: Công đoạn:
        B: AA / MA / LA ...
        """
        for i, v in enumerate(block[0]):
            if str(v).strip() == "Công đoạn:":
                cd_code = str(block.loc[i, 1]).strip().upper()
                return cd_code
        return None
  
    print("===== DUMP SHAPES PPT =====")
    for idx, slide in enumerate(prs.slides):
        print(f"\nSlide index {idx}:")
        for shape in slide.shapes:
            print("   ", repr(shape.name))

    for i, block in enumerate(blocks):
        cd_code = get_cong_doan_name(block)

        if not cd_code:
            print("⏭️ Bỏ qua block không có Công đoạn")
            continue

        print(f"📌 Công đoạn {cd_code}")

        if cd_code not in BLOCK_SLIDE_MAP:
            print(f"⚠️ Không có slide cho công đoạn {cd_code}")
            continue

        slide_idx = BLOCK_SLIDE_MAP[cd_code]
        title_text = f"TÍNH SẢN XUẤT – TỈ LỆ HOẠT ĐỘNG MÁY CÔNG ĐOẠN {cd_code}"
        
        update_title_in_ppt(
            prs,
            slide_idx,
            title_text,
            title_shape_name="title_block",
            font_name="Calibri",
            font_size_pt=20,
            bold=True,
            color_rgb=(0, 0, 128)   # xanh đậm
        )


        for cfg in BLOCK_CONFIG:
            df_daily, df_month = build_chart_df_block(block, cfg["rows"])

            img_path_local = LOCAL_OUT_DIR / f"{cfg['name']}_{i+1}.png"

            draw_combo_chart_block_excel_style(
                df_daily=df_daily,
                df_month=df_month,
                title=f"{cfg['title_prefix']} HOÀN THÀNH – CĐ {cd_code} – CA 1 + CA 2",
                ylabel="Số lượng",
                output_file=img_path_local,
                unit=cfg["unit"]
            )

            anchor = PPT_ANCHORS["block"][cfg["name"]]


            # ✅ COPY LÊN SERVER
            copy_to_server_safe(img_path_local, SERVER_OUT_DIR)

            replace_image_in_ppt(
                prs,
                anchor,
                img_path_local,
                slide_index=slide_idx
            )

    process_year_summary_combo(prs)
    # =================================================
    # FILE 3 – TABLE (G1:L16) → SLIDE 3
    # =================================================
    df_file3 = load_file3_G1_L16()

    # ✅ TIÊU ĐỀ TABLE (NẰM TRÊN)
    report_date = (datetime.now() - timedelta(days=1)).strftime("%d/%m/%Y")

    update_table_title_by_anchor(
        prs,
        slide_index=2,
        anchor_name="title_file3_table",
        title_text=f"BẢN VẼ HOÀN THÀNH, TRỄ HẸN VÀ PHỤ TẢI CÔNG ĐOẠN – {report_date}",
        font_name="Calibri",
        font_size_pt=20,
        bold=True,
        color_rgb=(0, 0, 139)   # xanh đậm
    )


    report_date = datetime.now().strftime("%Y-%m-%d")

    replace_table_by_anchor_centered(
        prs,
        slide_index=2,
        anchor_name="anchor_file3_table",
        df=df_file3,
        report_date=(datetime.now() - timedelta(days=1)).strftime("%Y-%m-%d")
    )
    process_costdown_noitac(prs)
    process_costdown_mt_tt_month(prs)
    
    update_title_in_ppt(
        prs,
        slide_index=18,                      # slide 19
        title_text="COSTDOWN NỘI TÁC – TỔNG HỢP NĂM 2026",
        title_shape_name="title_costdown",   # tên textbox title trong PPT
        font_name="Calibri",
        font_size_pt=20,
        bold=True,
        color_rgb=(0, 0, 128)                # xanh đậm
    )

    process_timer_utilization_chart(prs)
    update_title_in_ppt(
        prs,
        slide_index=3,                      # ✅ SLIDE 4
        title_text="TỈ LỆ HOẠT ĐỘNG MÁY TẠI CÁC CÔNG ĐOẠN TRONG THÁNG",
        title_shape_name="title_timer",     # ✅ đúng tên textbox title trong PPT
        font_name="Calibri",
        font_size_pt=20,
        bold=True,
        color_rgb=(0, 0, 128)               # xanh đậm
    )

    process_timer_block_charts(prs)

    # ==============================
    # SAVE PPT – SAFE MODE
    # ==============================
    file_view = PPT_FILE.parent / f"{PPT_FILE.stem}_view.pptx"
    file_final = PPT_FILE.parent / f"{PPT_FILE.stem}.pptx"

    def save_ppt_safe(prs, target: Path):
        tmp = target.with_suffix(".tmp.pptx")
        prs.save(tmp)
        os.replace(tmp, target)   # ✅ atomic replace – né lock / antivirus

    # ✅ Lưu bản xem
    save_ppt_safe(prs, file_view)
    print(f"✅ Đã lưu file view: {file_view.name}")

    # ✅ Lưu bản final (readonly + hidden)
    save_ppt_safe(prs, file_final)

    os.chmod(file_final, stat.S_IREAD)
    subprocess.run(
        f'attrib +H "{file_final}"',
        shell=True,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL
    )

    print(f"✅ Đã lưu file final (ẩn + readonly): {file_final.name}")
    print("🎉 HOÀN TẤT")


if __name__ == "__main__":
    main()