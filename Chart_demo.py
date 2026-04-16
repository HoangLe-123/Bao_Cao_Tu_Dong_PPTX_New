import pandas as pd
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pathlib import Path

# ======================================================
# CONFIG CHUNG
# ======================================================
EXCEL_FILE = r"Hằng ngày_data.xlsx"
EXCEL_SHEET = "Hằng ngày"

PPT_FILE = r"Bao_Cao_Hang_Ngay.pptx"
ROW_DATE = "Ngày làm việc"

OUT_DIR = Path("charts")
OUT_DIR.mkdir(exist_ok=True)

# ======================================================
# CẤU HÌNH CÁC BIỂU ĐỒ (CHỈ CẦN THÊM Ở ĐÂY)
# ======================================================
CHART_CONFIGS = [
    {
        "name": "ban_ve",
        "ppt_shape": "combo_ban_ve",
        "rows": {
            "tt": "THỰC TÍCH SỐ BẢN VẼ HOÀN THÀNH TRONG NGÀY",
            "mt": "MỤC TIÊU SỐ BẢN VẼ HOÀN THÀNH TRONG NGÀY",
            "lk": "MỤC TIÊU LŨY KẾ BẢN VẼ HT TRONG NGÀY ( CÒN LẠI)",
        },
        "title": "THỰC TÍCH SỐ BẢN VẼ HOÀN THÀNH THÁNG 04/2026",
        "ylabel": "Số bản vẽ",
        "legend": (
            "Thực tích bản vẽ",
            "Mục tiêu bản vẽ",
            "Mục tiêu lũy kế bản vẽ",
            "Thực tích bản vẽ1",
            "Mục tiêu bản vẽ2",
            "Mục tiêu lũy kế bản vẽ3",
            "Thực tích bản vẽ4",
            "Mục tiêu bản vẽ"
        )
    },
    # 🔹 Thêm biểu đồ khác (ví dụ PCS) ở đây
]

# ======================================================
# HÀM LOAD DATA GỐC
# ======================================================
def load_raw_excel():
    return pd.read_excel(EXCEL_FILE, sheet_name=EXCEL_SHEET, header=None)

# ======================================================
# HÀM BUILD DATA CHO 1 BIỂU ĐỒ
# ======================================================
def build_chart_df(df_raw, row_tt, row_mt, row_lk):
    i_date = df_raw.index[df_raw[0] == ROW_DATE][0]
    i_tt = df_raw.index[df_raw[0] == row_tt][0]
    i_mt = df_raw.index[df_raw[0] == row_mt][0]
    i_lk = df_raw.index[df_raw[0] == row_lk][0]

    dates = pd.to_datetime(df_raw.loc[i_date, 1:], errors="coerce")
    tt = df_raw.loc[i_tt, 1:].values
    mt = df_raw.loc[i_mt, 1:].values
    lk = df_raw.loc[i_lk, 1:].values

    df = pd.DataFrame({
        "Ngày": dates,
        "Thực tích": tt,
        "Mục tiêu": mt,
        "Lũy kế": lk
    }).dropna(subset=["Ngày"]).fillna(0)

    # Ẩn line khi = 0
    df["Mục tiêu"] = df["Mục tiêu"].replace(0, np.nan)
    df["Lũy kế"] = df["Lũy kế"].replace(0, np.nan)

    # Trục X dạng categorical
    df = df.reset_index(drop=True)
    df["X"] = df.index
    df["Label"] = df["Ngày"].dt.strftime("%d-%m")

    return df

# ======================================================
# HÀM VẼ COMBO CHART
# ======================================================
def draw_combo_chart(df, title, ylabel, legends, output_file):
    legend_bar = legends[0] if len(legends) > 0 else None
    legend_l1  = legends[1] if len(legends) > 1 else None
    legend_l2  = legends[2] if len(legends) > 2 else None

    plt.figure(figsize=(12, 5))

    # BAR
    bar_df = df[df["Thực tích"] > 0]
    bars = plt.bar(
        bar_df["X"],
        bar_df["Thực tích"],
        width=0.8,
        color="#9c1c7d",
        label=legend_bar
    )

    for bar in bars:
        h = bar.get_height()
        if h > 0:
            plt.text(
                bar.get_x() + bar.get_width()/2,
                h*1.01,
                f"{int(h)}",
                ha="center",
                va="bottom",
                fontsize=9
            )

    # LINE 1
    if legend_l1:
        plt.plot(
            df["X"], df["Mục tiêu"],
            color="green", marker="o",
            linewidth=2, label=legend_l1
        )

    # LINE 2
    if legend_l2:
        plt.plot(
            df["X"], df["Lũy kế"],
            color="red", marker="o",
            linewidth=2, label=legend_l2
        )

    # STYLE
    plt.title(title, fontsize=12)
    plt.ylabel(ylabel)
    plt.grid(axis="y", linestyle="--", alpha=0.4)

    max_ticks = 10
    step = max(1, len(df)//max_ticks)
    plt.xticks(df["X"].iloc[::step], df["Label"].iloc[::step])
    plt.tick_params(axis="x", labelsize=8)

    legend_count = len([l for l in [legend_bar, legend_l1, legend_l2] if l])

    plt.legend(
        loc="upper center",
        bbox_to_anchor=(0.5, -0.12 if legend_count <= 2 else -0.20),
        ncol=legend_count,
        frameon=False,
        fontsize=9
    )

    plt.subplots_adjust(
        left=0.06,
        right=0.995,
        top=0.90,
        bottom=0.24 if legend_count <= 2 else 0.32
    )

    plt.savefig(output_file, dpi=150)
    plt.close()
# ======================================================
# HÀM GHI ẢNH VÀO POWERPOINT
# ======================================================
def replace_image_in_ppt(shape_name, image_path):
    prs = Presentation(PPT_FILE)
    for slide in prs.slides:
        for shape in list(slide.shapes):
            if shape.name == shape_name:
                left, top = shape.left, shape.top
                width, height = shape.width, shape.height
                slide.shapes._spTree.remove(shape._element)
                slide.shapes.add_picture(
                    str(image_path),
                    left, top,
                    width=width, height=height
                )
    prs.save(PPT_FILE)

# ======================================================
# MAIN
# ======================================================
df_raw = load_raw_excel()

for cfg in CHART_CONFIGS:
    df_chart = build_chart_df(
        df_raw,
        cfg["rows"]["tt"],
        cfg["rows"]["mt"],
        cfg["rows"]["lk"]
    )

    img_path = OUT_DIR / f"{cfg['name']}.png"

    draw_combo_chart(
        df=df_chart,
        title=cfg["title"],
        ylabel=cfg["ylabel"],
        legends=cfg["legend"],
        output_file=img_path
    )

    replace_image_in_ppt(cfg["ppt_shape"], img_path)
    print(f"✅ Đã cập nhật biểu đồ: {cfg['title']}")

print("🎉 HOÀN TẤT TOÀN BỘ BIỂU ĐỒ")